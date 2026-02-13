"""
PowerPoint Parser
Reads PowerPoint files and extracts slides, notes, and metadata
"""

import os
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed")
    print("Install with: pip install python-pptx")
    raise


class PowerPointParser:
    """Parse PowerPoint presentations and extract content"""

    def __init__(self, pptx_file):
        """
        Initialize parser with PowerPoint file

        Args:
            pptx_file: Path to .pptx file
        """
        self.pptx_file = pptx_file
        self.presentation = None
        self.slides_data = []

    def parse(self, preview_count=None):
        """
        Parse PowerPoint and extract all slide data

        Args:
            preview_count: Only parse first N slides (None = all)

        Returns:
            List of slide dictionaries with:
            - slide_number: int
            - title: str
            - content: str (bullet points, text)
            - notes: str (speaker notes)
            - image_path: str (exported slide image)
            - broll_markers: list (detected [SCREEN:file] markers)
            - pause_markers: list (detected [PAUSE:seconds] markers)
        """
        print(f"   Loading: {self.pptx_file}")

        try:
            self.presentation = Presentation(self.pptx_file)
        except Exception as e:
            raise Exception(f"Failed to load PowerPoint: {e}")

        total_slides = len(self.presentation.slides)
        if preview_count:
            slides_to_parse = min(preview_count, total_slides)
            print(f"   Preview mode: {slides_to_parse}/{total_slides} slides")
        else:
            slides_to_parse = total_slides

        print(f"   Parsing {slides_to_parse} slides...")

        for i, slide in enumerate(self.presentation.slides):
            if preview_count and i >= preview_count:
                break

            slide_data = self._parse_slide(slide, i + 1)
            self.slides_data.append(slide_data)

        return self.slides_data

    def _parse_slide(self, slide, slide_number):
        """Parse individual slide"""

        # Extract title
        title = self._extract_title(slide)

        # Extract text content
        content = self._extract_content(slide)

        # Extract speaker notes
        notes = self._extract_notes(slide)

        # Parse special markers in notes
        broll_markers = self._extract_broll_markers(notes)
        pause_markers = self._extract_pause_markers(notes)
        speed_markers = self._extract_speed_markers(notes)

        # Clean notes (remove markers for voice generation)
        clean_notes = self._clean_notes(notes)

        return {
            'slide_number': slide_number,
            'title': title,
            'content': content,
            'notes': clean_notes,
            'notes_raw': notes,  # Keep original with markers
            'broll_markers': broll_markers,
            'pause_markers': pause_markers,
            'speed_markers': speed_markers,
            'slide_object': slide  # Keep reference to slide
        }

    def _extract_title(self, slide):
        """Extract slide title"""
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.has_text_frame:
                    if hasattr(shape, 'placeholder_format'):
                        # This is likely a title placeholder
                        if shape.text.strip():
                            return shape.text.strip()

            # Fallback: first text box
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    return shape.text.strip()

        return f"Slide {slide.slide_id}"

    def _extract_content(self, slide):
        """Extract all text content from slide"""
        content_parts = []

        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        content_parts.append(text)

        return "\n".join(content_parts)

    def _extract_notes(self, slide):
        """Extract speaker notes"""
        try:
            if hasattr(slide, 'notes_slide'):
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame'):
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    return notes_text
        except Exception:
            pass

        return ""

    def _extract_broll_markers(self, notes):
        """
        Extract [SCREEN:filename] markers from notes

        Examples:
            [SCREEN:demo.mp4]
            [SCREEN:demos/intro.mp4]
            [BROLL:cutaway.mp4]
        """
        markers = []
        patterns = [
            r'\[SCREEN:([^\]]+)\]',
            r'\[BROLL:([^\]]+)\]'
        ]

        for pattern in patterns:
            matches = re.finditer(pattern, notes, re.IGNORECASE)
            for match in matches:
                filename = match.group(1).strip()
                markers.append({
                    'type': 'screen',
                    'file': filename,
                    'position': match.start()
                })

        return markers

    def _extract_pause_markers(self, notes):
        """
        Extract [PAUSE:seconds] markers from notes

        Examples:
            [PAUSE:2]
            [PAUSE:1.5]
        """
        markers = []
        pattern = r'\[PAUSE:(\d+\.?\d*)\]'

        matches = re.finditer(pattern, notes, re.IGNORECASE)
        for match in matches:
            duration = float(match.group(1))
            markers.append({
                'type': 'pause',
                'duration': duration,
                'position': match.start()
            })

        return markers

    def _extract_speed_markers(self, notes):
        """
        Extract [SLOW] or [FAST] markers from notes

        Examples:
            [SLOW]
            [FAST]
            [SPEED:0.8]
        """
        markers = []

        # Named speed markers
        for match in re.finditer(r'\[(SLOW|FAST)\]', notes, re.IGNORECASE):
            speed = 0.85 if match.group(1).upper() == 'SLOW' else 1.15
            markers.append({
                'type': 'speed',
                'factor': speed,
                'position': match.start()
            })

        # Numeric speed markers
        for match in re.finditer(r'\[SPEED:(\d+\.?\d*)\]', notes, re.IGNORECASE):
            speed = float(match.group(1))
            markers.append({
                'type': 'speed',
                'factor': speed,
                'position': match.start()
            })

        return markers

    def _clean_notes(self, notes):
        """Remove all special markers from notes for voice generation"""
        if not notes:
            return ""

        # Remove all markers
        clean = notes
        patterns = [
            r'\[SCREEN:[^\]]+\]',
            r'\[BROLL:[^\]]+\]',
            r'\[PAUSE:\d+\.?\d*\]',
            r'\[SLOW\]',
            r'\[FAST\]',
            r'\[SPEED:\d+\.?\d*\]',
            r'\[HIGHLIGHT:[^\]]+\]',
            r'\[ZOOM:[^\]]+\]'
        ]

        for pattern in patterns:
            clean = re.sub(pattern, '', clean, flags=re.IGNORECASE)

        # Clean up extra whitespace
        clean = re.sub(r'\s+', ' ', clean).strip()

        return clean

    def export_slides_as_images(self, output_dir):
        """
        Export each slide as an image file

        Args:
            output_dir: Directory to save images

        Returns:
            List of image file paths
        """
        os.makedirs(output_dir, exist_ok=True)

        # Note: python-pptx doesn't support direct image export
        # We need to use PowerPoint COM API (Windows) or conversion tool
        print("   ⚠️  Slide image export requires PowerPoint COM API")
        print("   Falling back to placeholder method...")

        # This would need PowerPoint installed on Windows
        # For now, return placeholder paths
        image_paths = []
        for i, slide_data in enumerate(self.slides_data, 1):
            image_path = f"{output_dir}/slide_{i:03d}.png"
            image_paths.append(image_path)

        return image_paths


def test_parser():
    """Test the parser with a sample file"""
    import sys

    if len(sys.argv) < 2:
        print("Usage: python ppt_parser.py presentation.pptx")
        return

    pptx_file = sys.argv[1]

    print("=" * 70)
    print("PowerPoint Parser Test")
    print("=" * 70)

    parser = PowerPointParser(pptx_file)
    slides = parser.parse()

    print(f"\nFound {len(slides)} slides:\n")

    for slide in slides:
        print(f"Slide {slide['slide_number']}: {slide['title']}")
        print(f"  Content: {slide['content'][:100]}...")
        print(f"  Notes: {slide['notes'][:100] if slide['notes'] else '(no notes)'}...")

        if slide['broll_markers']:
            print(f"  B-roll: {[m['file'] for m in slide['broll_markers']]}")
        if slide['pause_markers']:
            print(f"  Pauses: {[m['duration'] for m in slide['pause_markers']]}")

        print()


if __name__ == "__main__":
    test_parser()
