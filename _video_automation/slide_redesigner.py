"""
Automated Slide Redesigner - Data Engineering Focus
Transforms text-heavy slides into visual, engaging presentations
Specialized for Snowflake, Databricks, and Data Engineering topics
"""

import sys
import os
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

class SlideRedesigner:
    """Redesigns PowerPoint slides for maximum visual impact"""

    def __init__(self, input_pptx, output_pptx=None):
        self.input_pptx = input_pptx
        self.output_pptx = output_pptx or self._generate_output_name()
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(9)

        # Color schemes for data engineering
        self.colors = {
            'snowflake_blue': RGBColor(41, 179, 209),  # Snowflake brand color
            'databricks_red': RGBColor(255, 58, 33),    # Databricks brand color
            'dark_bg': RGBColor(26, 35, 46),            # Dark professional background
            'light_text': RGBColor(255, 255, 255),      # White text
            'accent_teal': RGBColor(0, 188, 212),       # Accent color
            'accent_orange': RGBColor(255, 152, 0),     # Accent color
        }

    def _generate_output_name(self):
        """Generate output filename"""
        path = Path(self.input_pptx)
        return str(path.parent / f"{path.stem}_redesigned.pptx")

    def analyze_slide_type(self, slide_content):
        """Determine what type of slide this is based on content"""
        content_lower = slide_content.lower()

        # Check for comparison keywords
        if any(word in content_lower for word in ['vs', 'versus', 'comparison', 'compare', 'difference']):
            return 'comparison'

        # Check for architecture keywords
        if any(word in content_lower for word in ['architecture', 'pipeline', 'flow', 'design', 'infrastructure', 'system']):
            return 'architecture'

        # Check for feature/list keywords
        if any(word in content_lower for word in ['feature', 'benefit', 'advantage', 'capability', 'key points']):
            return 'feature_list'

        # Check for data/metrics keywords
        if any(word in content_lower for word in ['performance', 'metrics', 'benchmark', 'speed', 'cost', 'pricing']):
            return 'data_metrics'

        # Check for technical/code keywords
        if any(word in content_lower for word in ['sql', 'code', 'query', 'example', 'syntax']):
            return 'technical'

        # Default to simple visual
        return 'simple_visual'

    def extract_key_points(self, slide_text, max_points=4):
        """Extract key bullet points from slide text"""
        # Split by bullet points or newlines
        lines = re.split(r'[\n•\-]', slide_text)
        points = [line.strip() for line in lines if line.strip() and len(line.strip()) > 10]

        # Return top N points
        return points[:max_points]

    def create_title_slide(self, title, subtitle=""):
        """Create modern title slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Dark gradient background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['dark_bg']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(14), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['light_text']

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5),
                Inches(14), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = self.colors['accent_teal']

        # Accent line
        line = slide.shapes.add_shape(
            1,  # Line shape
            Inches(6), Inches(5),
            Inches(4), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['snowflake_blue']

        return slide

    def create_comparison_slide(self, title, content, notes=""):
        """Create comparison slide (Snowflake vs Databricks style)"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['dark_bg']

        # Title
        self._add_slide_title(slide, title)

        # Create two-column comparison
        col_width = Inches(6)
        col_height = Inches(5)

        # Snowflake column (left)
        left_box = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1.5), Inches(2),
            col_width, col_height
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
        left_box.line.color.rgb = self.colors['snowflake_blue']
        left_box.line.width = Pt(3)

        # Databricks column (right)
        right_box = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(8.5), Inches(2),
            col_width, col_height
        )
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
        right_box.line.color.rgb = self.colors['databricks_red']
        right_box.line.width = Pt(3)

        # Extract comparison points from content
        points = self.extract_key_points(content, max_points=6)

        # Add headers
        left_header = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.2),
            col_width, Inches(0.6)
        )
        left_header_frame = left_header.text_frame
        left_header_frame.text = "SNOWFLAKE"
        left_header_para = left_header_frame.paragraphs[0]
        left_header_para.alignment = PP_ALIGN.CENTER
        left_header_para.font.size = Pt(24)
        left_header_para.font.bold = True
        left_header_para.font.color.rgb = self.colors['snowflake_blue']

        right_header = slide.shapes.add_textbox(
            Inches(8.5), Inches(2.2),
            col_width, Inches(0.6)
        )
        right_header_frame = right_header.text_frame
        right_header_frame.text = "DATABRICKS"
        right_header_para = right_header_frame.paragraphs[0]
        right_header_para.alignment = PP_ALIGN.CENTER
        right_header_para.font.size = Pt(24)
        right_header_para.font.bold = True
        right_header_para.font.color.rgb = self.colors['databricks_red']

        # Add comparison points (split between columns)
        mid_point = len(points) // 2

        # Left column points
        left_text = slide.shapes.add_textbox(
            Inches(2), Inches(3.2),
            Inches(5), Inches(3.5)
        )
        left_frame = left_text.text_frame
        for point in points[:mid_point]:
            p = left_frame.add_paragraph()
            p.text = f"• {point[:80]}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['light_text']
            p.space_after = Pt(12)

        # Right column points
        right_text = slide.shapes.add_textbox(
            Inches(9), Inches(3.2),
            Inches(5), Inches(3.5)
        )
        right_frame = right_text.text_frame
        for point in points[mid_point:]:
            p = right_frame.add_paragraph()
            p.text = f"• {point[:80]}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['light_text']
            p.space_after = Pt(12)

        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def create_architecture_slide(self, title, content, notes=""):
        """Create architecture diagram slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['dark_bg']

        # Title
        self._add_slide_title(slide, title)

        # Architecture placeholder (simplified version)
        # In a full implementation, this would generate actual diagrams

        # Main architecture box
        arch_box = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(2), Inches(2.5),
            Inches(12), Inches(4.5)
        )
        arch_box.fill.solid()
        arch_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
        arch_box.line.color.rgb = self.colors['accent_teal']
        arch_box.line.width = Pt(2)

        # Architecture label
        arch_label = slide.shapes.add_textbox(
            Inches(2.5), Inches(3),
            Inches(11), Inches(1)
        )
        arch_frame = arch_label.text_frame
        arch_frame.text = "📊 Architecture Diagram"
        arch_para = arch_frame.paragraphs[0]
        arch_para.alignment = PP_ALIGN.CENTER
        arch_para.font.size = Pt(32)
        arch_para.font.color.rgb = self.colors['accent_teal']

        # Key components (extracted from content)
        points = self.extract_key_points(content, max_points=4)

        # Add component boxes
        y_start = 4.5
        for i, point in enumerate(points):
            comp_box = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(3 + (i * 2.5)), Inches(y_start),
                Inches(2), Inches(0.8)
            )
            comp_box.fill.solid()
            comp_box.fill.fore_color.rgb = self.colors['snowflake_blue']

            comp_text = slide.shapes.add_textbox(
                Inches(3 + (i * 2.5)), Inches(y_start + 0.1),
                Inches(2), Inches(0.6)
            )
            comp_frame = comp_text.text_frame
            comp_frame.text = point[:20]
            comp_para = comp_frame.paragraphs[0]
            comp_para.alignment = PP_ALIGN.CENTER
            comp_para.font.size = Pt(12)
            comp_para.font.bold = True
            comp_para.font.color.rgb = self.colors['light_text']

        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def create_feature_list_slide(self, title, content, notes=""):
        """Create feature list slide with icons"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['dark_bg']

        # Title
        self._add_slide_title(slide, title)

        # Extract features
        features = self.extract_key_points(content, max_points=6)

        # Icons for features (using emojis as placeholders)
        icons = ['✓', '⚡', '🔒', '📊', '🚀', '💡']

        # Create 2x3 grid of features
        rows = 2
        cols = 3
        box_width = Inches(4)
        box_height = Inches(2)
        spacing_x = Inches(0.5)
        spacing_y = Inches(0.3)
        start_x = Inches(1.5)
        start_y = Inches(2.5)

        for i, feature in enumerate(features):
            row = i // cols
            col = i % cols

            x = start_x + col * (box_width + spacing_x)
            y = start_y + row * (box_height + spacing_y)

            # Feature box
            feat_box = slide.shapes.add_shape(
                1,  # Rectangle
                x, y,
                box_width, box_height
            )
            feat_box.fill.solid()
            feat_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
            feat_box.line.color.rgb = self.colors['accent_teal']
            feat_box.line.width = Pt(2)

            # Icon
            icon_text = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.2),
                Inches(0.8), Inches(0.8)
            )
            icon_frame = icon_text.text_frame
            icon_frame.text = icons[i % len(icons)]
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(36)
            icon_para.font.color.rgb = self.colors['accent_teal']

            # Feature text
            feat_text = slide.shapes.add_textbox(
                x + Inches(1.2), y + Inches(0.3),
                Inches(2.5), Inches(1.5)
            )
            feat_frame = feat_text.text_frame
            feat_frame.text = feature[:100]
            feat_frame.word_wrap = True
            feat_para = feat_frame.paragraphs[0]
            feat_para.font.size = Pt(14)
            feat_para.font.color.rgb = self.colors['light_text']

        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def create_simple_visual_slide(self, title, content, notes=""):
        """Create simple visual slide with minimal text"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['dark_bg']

        # Title
        self._add_slide_title(slide, title)

        # Extract key message (just 1-2 sentences)
        points = self.extract_key_points(content, max_points=2)

        # Large central text box
        main_text = slide.shapes.add_textbox(
            Inches(2), Inches(3.5),
            Inches(12), Inches(3)
        )
        main_frame = main_text.text_frame
        main_frame.text = "\n\n".join(points)
        main_frame.word_wrap = True

        for para in main_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(32)
            para.font.color.rgb = self.colors['light_text']
            para.space_after = Pt(20)

        # Accent decoration
        accent = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(6), Inches(6.8),
            Inches(4), Inches(0.1)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.colors['accent_teal']

        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def _add_slide_title(self, slide, title):
        """Add consistent title to slide"""
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5),
            Inches(14), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['light_text']

        # Title underline
        underline = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1), Inches(1.45),
            Inches(3), Inches(0.05)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self.colors['snowflake_blue']

    def redesign_presentation(self, input_prs):
        """Main redesign logic"""
        print(f"\nAnalyzing presentation: {self.input_pptx}")
        print("-" * 70)

        slide_count = 0

        for slide_num, slide in enumerate(input_prs.slides, 1):
            # Extract slide content
            title = ""
            content = ""
            notes = ""

            # Get title
            if slide.shapes.title:
                title = slide.shapes.title.text

            # Get content from all text boxes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"

            # Get speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text

            # Skip if empty
            if not title and not content.strip():
                continue

            print(f"[{slide_num}] {title[:60]}...")

            # Determine slide type
            slide_type = self.analyze_slide_type(content)
            print(f"     Type: {slide_type}")

            # Create appropriate slide
            if slide_num == 1:
                # Title slide
                subtitle = notes[:100] if notes else ""
                self.create_title_slide(title, subtitle)
            elif slide_type == 'comparison':
                self.create_comparison_slide(title, content, notes)
            elif slide_type == 'architecture':
                self.create_architecture_slide(title, content, notes)
            elif slide_type == 'feature_list':
                self.create_feature_list_slide(title, content, notes)
            else:
                self.create_simple_visual_slide(title, content, notes)

            slide_count += 1

        return slide_count

    def save(self):
        """Save redesigned presentation"""
        self.prs.save(self.output_pptx)
        print(f"\n✓ Saved: {self.output_pptx}")


def main():
    parser = argparse.ArgumentParser(
        description="Redesign PowerPoint slides for Data Engineering topics",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Basic redesign
  python slide_redesigner.py presentation.pptx

  # Specify output file
  python slide_redesigner.py input.pptx --output redesigned.pptx

  # Then create video
  python video_creator.py redesigned.pptx --audio-only
        """
    )

    parser.add_argument('presentation',
                       help='Input PowerPoint file (.pptx)')

    parser.add_argument('--output', '-o',
                       help='Output PowerPoint file (default: [name]_redesigned.pptx)')

    args = parser.parse_args()

    # Validate input file
    if not os.path.exists(args.presentation):
        print(f"ERROR: Presentation file not found: {args.presentation}")
        return 1

    print("=" * 70)
    print("SLIDE REDESIGNER - Data Engineering Visual Transformation")
    print("=" * 70)

    try:
        # Load original presentation
        print(f"\nLoading: {args.presentation}")
        input_prs = Presentation(args.presentation)

        # Create redesigner
        redesigner = SlideRedesigner(args.presentation, args.output)

        # Redesign slides
        slide_count = redesigner.redesign_presentation(input_prs)

        # Save
        redesigner.save()

        # Summary
        print("\n" + "=" * 70)
        print("SUCCESS! Slide redesign complete!")
        print("=" * 70)
        print(f"\nOriginal: {args.presentation}")
        print(f"Redesigned: {redesigner.output_pptx}")
        print(f"Slides processed: {slide_count}")

        print("\nNext steps:")
        print(f"   1. Review: {redesigner.output_pptx}")
        print(f"   2. Generate audio: python video_creator.py {redesigner.output_pptx} --audio-only")
        print(f"   3. Create video: python video_creator.py {redesigner.output_pptx} --use-existing-audio")

        return 0

    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())
