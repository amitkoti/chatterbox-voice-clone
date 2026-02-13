"""
Snowbrix Slide Composer
Professional cream + green slide design system
"""

from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from brand_colors_snowbrix import SnowbrixColors
from PIL import Image as PILImage
from lxml.etree import Element
import io


class SnowbrixSlideComposer:
    """Slide composer with Snowbrix professional design"""

    def __init__(self, output_path: str, include_logo: bool = True):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.include_logo = include_logo

        # Get Snowbrix colors
        scheme = SnowbrixColors.get_scheme()
        self.colors = scheme
        print(f"Using: {scheme['name']}")

    def _remove_bullets(self, paragraph):
        """Remove bullet formatting from paragraph"""
        try:
            pPr = paragraph._element.get_or_add_pPr()
            for child in list(pPr):
                if 'bu' in child.tag.lower():
                    pPr.remove(child)
            buNone = Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
            pPr.insert(0, buNone)
        except:
            pass

    def _add_logo_to_slide(self, slide):
        """Add Snowbrix text logo to slide (bottom right corner)"""
        if not self.include_logo:
            return

        # Logo position (bottom right corner - equal margins)
        logo_x = 13.7
        logo_y = 8.3
        font_size = Pt(22)

        # Use single textbox with colored runs for perfect spacing
        logo_box = slide.shapes.add_textbox(
            Inches(logo_x), Inches(logo_y),
            Inches(2.5), Inches(0.5)
        )

        text_frame = logo_box.text_frame
        text_frame.clear()

        paragraph = text_frame.paragraphs[0]

        # Add "SNOW" in forest green
        run_snow = paragraph.add_run()
        run_snow.text = "SNOW"
        run_snow.font.size = font_size
        run_snow.font.bold = True
        run_snow.font.color.rgb = self.colors['color_primary']

        # Add "BRIX" in terracotta (no space)
        run_brix = paragraph.add_run()
        run_brix.text = "BRIX"
        run_brix.font.size = font_size
        run_brix.font.bold = True
        run_brix.font.color.rgb = self.colors['accent_terracotta']

        # Add "AI" in charcoal (no space)
        run_ai = paragraph.add_run()
        run_ai.text = "AI"
        run_ai.font.size = font_size
        run_ai.font.bold = True
        run_ai.font.color.rgb = self.colors['text_secondary']

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Title slide - Snowbrix style"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Large subtle green shape (top right)
        shape = slide.shapes.add_shape(
            1, Inches(12), Inches(-0.5),
            Inches(5), Inches(5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['color_primary_light']
        shape.fill.transparency = 0.5
        shape.line.fill.background()
        shape.rotation = 20

        # Thick Forest Green accent bar (left)
        bar = slide.shapes.add_shape(
            1, Inches(1.5), Inches(3.5),
            Inches(0.3), Inches(3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['color_primary']
        bar.line.fill.background()

        # Title in FOREST GREEN (#2D5F3F)
        title_box = slide.shapes.add_textbox(
            Inches(2.2), Inches(3.8),
            Inches(12), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']  # Forest Green

        # Subtitle in SAGE GREEN (#6B9080)
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(2.2), Inches(6),
                Inches(12), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.LEFT
            sub_para.font.size = Pt(36)
            sub_para.font.color.rgb = self.colors['color_secondary']  # Sage Green

        # Add logo
        self._add_logo_to_slide(slide)

        return slide

    def add_content_slide(
        self,
        title: str,
        key_points: List[str] = None,
        notes: str = ""
    ):
        """Content slide - Snowbrix style with strong contrast"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Light cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_secondary']  # #FAF5F0

        # === TITLE SECTION ===
        # Subtle green background for title area
        title_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(16), Inches(2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['bg_accent']  # Light green #E8F3ED
        title_bg.line.fill.background()

        # Thick accent bar (Sage Green)
        bar = slide.shapes.add_shape(
            1, Inches(0.8), Inches(0.5),
            Inches(0.2), Inches(1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['color_secondary']  # Sage Green
        bar.line.fill.background()

        # Title in FOREST GREEN (#2D5F3F) - DARK & BOLD
        title_box = slide.shapes.add_textbox(
            Inches(1.3), Inches(0.6),
            Inches(14), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']  # Forest Green - STANDS OUT!

        # === KEY POINTS with numbered circles ===
        if key_points:
            # Adjust spacing based on number of points
            num_points = min(len(key_points), 5)  # Max 5 points
            if num_points <= 3:
                start_y, spacing = 3.0, 1.6
            elif num_points == 4:
                start_y, spacing = 2.8, 1.3
            else:  # 5 points
                start_y, spacing = 2.6, 1.15

            for i, point in enumerate(key_points[:5]):
                y = start_y + (i * spacing)

                # Circle with SAGE GREEN (#6B9080)
                circle = slide.shapes.add_shape(
                    1, Inches(1.5), Inches(y),
                    Inches(0.7), Inches(0.7)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']  # Sage Green circle
                circle.line.fill.background()

                # Number (cream text on green)
                num_box = slide.shapes.add_textbox(
                    Inches(1.5), Inches(y),
                    Inches(0.7), Inches(0.7)
                )
                num_frame = num_box.text_frame
                num_frame.clear()  # Clear any default formatting
                num_para = num_frame.paragraphs[0]
                num_para.text = str(i + 1)  # Set text via paragraph, not frame
                num_para.alignment = PP_ALIGN.CENTER
                num_para.font.size = Pt(36)
                num_para.font.bold = True
                num_para.font.color.rgb = self.colors['text_on_green']  # Cream on green

                # Point text in WARM GRAY (#4A4A4A) - READABLE & DARK
                point_box = slide.shapes.add_textbox(
                    Inches(2.5), Inches(y + 0.05),
                    Inches(12.5), Inches(0.9)
                )
                point_frame = point_box.text_frame
                point_frame.clear()  # Clear any default formatting
                point_frame.word_wrap = True
                point_para = point_frame.paragraphs[0]
                point_para.text = point[:120]  # Set text via paragraph, not frame
                point_para.font.size = Pt(30)
                point_para.font.bold = True
                point_para.font.color.rgb = self.colors['text_secondary']  # Warm Gray - DARK!

                # Explicitly remove bullet formatting at XML level
                from lxml.etree import Element
                pPr = point_para._element.get_or_add_pPr()
                # Remove any existing bullet elements
                for child in list(pPr):
                    if 'bu' in child.tag.lower():
                        pPr.remove(child)
                # Add buNone element to explicitly disable bullets
                buNone = Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                pPr.insert(0, buNone)

        # Add notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        # Add logo
        self._add_logo_to_slide(slide)

        return slide

    def add_emphasis_slide(self, title: str, quote: str = ""):
        """Emphasis/Quote slide with light green background"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Light green background (#E8F3ED)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_accent']

        # Large quote or emphasis text (centered, with margins)
        text_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5),
            Inches(13), Inches(4)
        )
        text_frame = text_box.text_frame
        text_frame.text = quote or title
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.alignment = PP_ALIGN.CENTER
        text_para.font.size = Pt(44)
        text_para.font.bold = True
        text_para.font.color.rgb = self.colors['color_primary']  # Forest Green

        # Add logo
        self._add_logo_to_slide(slide)

        return slide

    def add_two_column_slide(
        self,
        title: str,
        left_content: List[str] = None,
        right_content: List[str] = None,
        left_title: str = "",
        right_title: str = "",
        notes: str = ""
    ):
        """Two-column layout for comparisons"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_secondary']

        # Title section with green background
        title_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(16), Inches(1.5)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['bg_accent']
        title_bg.line.fill.background()

        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5),
            Inches(14), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # LEFT COLUMN
        col_y = 2.2

        # Left column title (at left edge)
        if left_title:
            left_title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(col_y),  # At left edge, not indented
                Inches(6.5), Inches(0.5)
            )
            ltf = left_title_box.text_frame
            ltf.text = left_title
            ltp = ltf.paragraphs[0]
            ltp.font.size = Pt(32)
            ltp.font.bold = True
            ltp.font.color.rgb = self.colors['color_secondary']

        # Left content (bullets indented under heading)
        if left_content:
            content_y = col_y + (0.7 if left_title else 0)
            for i, item in enumerate(left_content[:5]):
                y_pos = content_y + i * 1.1  # Increased spacing from 0.8 to 1.1

                # Green bullet circle (indented under heading)
                circle = slide.shapes.add_shape(
                    1, Inches(1.0), Inches(y_pos + 0.1),
                    Inches(0.2), Inches(0.2)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Item text (indented further right, no text bullet)
                item_box = slide.shapes.add_textbox(
                    Inches(1.4), Inches(y_pos),  # Indented 0.6" from heading
                    Inches(6.0), Inches(0.9)  # Increased height
                )
                itf = item_box.text_frame
                itf.clear()
                itf.text = item  # No "• " prefix
                itf.word_wrap = True
                itp = itf.paragraphs[0]
                itp.font.size = Pt(24)  # Increased from 18 to 24
                itp.font.color.rgb = self.colors['text_secondary']
                self._remove_bullets(itp)

        # RIGHT COLUMN
        col_y = 2.2

        # Right column title (at column edge)
        if right_title:
            right_title_box = slide.shapes.add_textbox(
                Inches(8.5), Inches(col_y),  # At column edge
                Inches(6.5), Inches(0.5)
            )
            rtf = right_title_box.text_frame
            rtf.text = right_title
            rtp = rtf.paragraphs[0]
            rtp.font.size = Pt(32)
            rtp.font.bold = True
            rtp.font.color.rgb = self.colors['color_secondary']

        # Right content (bullets indented under heading)
        if right_content:
            content_y = col_y + (0.7 if right_title else 0)
            for i, item in enumerate(right_content[:5]):
                y_pos = content_y + i * 1.1  # Increased spacing from 0.8 to 1.1

                # Green bullet circle (indented under heading)
                circle = slide.shapes.add_shape(
                    1, Inches(8.7), Inches(y_pos + 0.1),
                    Inches(0.2), Inches(0.2)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Item text (indented further right, no text bullet)
                item_box = slide.shapes.add_textbox(
                    Inches(9.1), Inches(y_pos),  # Indented 0.6" from heading
                    Inches(6.0), Inches(0.9)  # Increased height
                )
                itf = item_box.text_frame
                itf.clear()
                itf.text = item  # No "• " prefix
                itf.word_wrap = True
                itp = itf.paragraphs[0]
                itp.font.size = Pt(24)  # Increased from 18 to 24
                itp.font.color.rgb = self.colors['text_secondary']
                self._remove_bullets(itp)

        # Add notes and logo
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        return slide

    def add_section_divider(self, section_title: str, subtitle: str = ""):
        """Section divider slide with large text"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Forest green background for impact
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['color_primary']

        # Large centered title (white text on green)
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5),
            Inches(12), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(64)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_on_green']

        # Subtitle if provided
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(3), Inches(5.2),
                Inches(10), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.CENTER
            sub_para.font.size = Pt(28)
            sub_para.font.color.rgb = self.colors['bg_primary']

        # Logo (white version for dark background)
        # Skip logo on section dividers for cleaner look

        return slide

    def add_image_text_split_slide(
        self,
        title: str,
        image_path: str = None,
        content_points: List[str] = None,
        image_on_left: bool = True,
        notes: str = ""
    ):
        """Image on one side, text on the other - very common layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Title bar with green accent
        title_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(16), Inches(1.3)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['bg_accent']
        title_bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(14), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Calculate positions based on image side
        if image_on_left:
            image_x, image_y, image_w, image_h = 0.8, 2.0, 7.0, 5.5
            text_x, text_y, text_w = 8.2, 2.2, 7.0
        else:
            text_x, text_y, text_w = 0.8, 2.2, 7.0
            image_x, image_y, image_w, image_h = 8.2, 2.0, 7.0, 5.5

        # Add image if provided
        if image_path and Path(image_path).exists():
            try:
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(image_x),
                    Inches(image_y),
                    width=Inches(image_w),
                    height=Inches(image_h)
                )
            except Exception as e:
                print(f"   Warning: Could not add image: {e}")
        else:
            # Placeholder box for image
            placeholder = slide.shapes.add_shape(
                1, Inches(image_x), Inches(image_y),
                Inches(image_w), Inches(image_h)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = self.colors['bg_accent']
            placeholder.line.color.rgb = self.colors['color_secondary']
            placeholder.line.width = Pt(2)

            # Placeholder text
            placeholder_text = slide.shapes.add_textbox(
                Inches(image_x + 2), Inches(image_y + 2.5),
                Inches(3), Inches(0.6)
            )
            ptf = placeholder_text.text_frame
            ptf.text = "[Image Placeholder]"
            ptp = ptf.paragraphs[0]
            ptp.alignment = PP_ALIGN.CENTER
            ptp.font.size = Pt(24)
            ptp.font.color.rgb = self.colors['color_secondary']

        # Add content points
        if content_points:
            for i, point in enumerate(content_points[:5]):
                y = text_y + (i * 1.0)

                # Bullet circle
                circle = slide.shapes.add_shape(
                    1, Inches(text_x), Inches(y),
                    Inches(0.35), Inches(0.35)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Point text
                point_box = slide.shapes.add_textbox(
                    Inches(text_x + 0.5), Inches(y - 0.05),
                    Inches(text_w - 0.6), Inches(0.8)
                )
                point_frame = point_box.text_frame
                point_frame.clear()
                point_frame.word_wrap = True
                point_para = point_frame.paragraphs[0]
                point_para.text = point
                point_para.font.size = Pt(24)
                point_para.font.color.rgb = self.colors['text_secondary']

                # Remove bullets
                from lxml.etree import Element
                pPr = point_para._element.get_or_add_pPr()
                for child in list(pPr):
                    if 'bu' in child.tag.lower():
                        pPr.remove(child)
                buNone = Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                pPr.insert(0, buNone)

        # Add notes and logo
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        return slide

    def add_thank_you_slide(self, message: str = "Thank You", contact_info: str = ""):
        """Thank you closing slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Light green background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_accent']

        # Large "Thank You" text
        thanks_box = slide.shapes.add_textbox(
            Inches(2), Inches(3),
            Inches(12), Inches(2)
        )
        thanks_frame = thanks_box.text_frame
        thanks_frame.text = message
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.alignment = PP_ALIGN.CENTER
        thanks_para.font.size = Pt(72)
        thanks_para.font.bold = True
        thanks_para.font.color.rgb = self.colors['color_primary']

        # Contact info
        if contact_info:
            contact_box = slide.shapes.add_textbox(
                Inches(3), Inches(5.5),
                Inches(10), Inches(1)
            )
            contact_frame = contact_box.text_frame
            contact_frame.text = contact_info
            contact_para = contact_frame.paragraphs[0]
            contact_para.alignment = PP_ALIGN.CENTER
            contact_para.font.size = Pt(24)
            contact_para.font.color.rgb = self.colors['color_secondary']

        self._add_logo_to_slide(slide)
        return slide

    def add_slide_with_image(
        self,
        title: str,
        image_path=None,
        slide_type: str = "content",
        notes: str = "",
        key_points: list = None
    ):
        """Compatibility method for slide redesigner integration"""
        if slide_type == "emphasis" or slide_type == "quote":
            return self.add_emphasis_slide(title, notes)
        else:
            # Standard content slide
            return self.add_content_slide(title, key_points or [], notes)

    def save(self):
        """Save presentation"""
        self.prs.save(self.output_path)
        print(f"[OK] Saved: {self.output_path}")


# Demo
if __name__ == "__main__":
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    composer = SnowbrixSlideComposer(f"_Snowbrix_Sample_{timestamp}.pptx")

    composer.add_title_slide(
        "Data Engineering Mastery",
        "Transform Data into Strategic Advantage"
    )

    composer.add_content_slide(
        "Why Snowflake Wins",
        key_points=[
            "10x faster queries than traditional data warehouses",
            "Zero infrastructure management - fully automated scaling",
            "Instant data sharing across teams and organizations"
        ]
    )

    composer.add_content_slide(
        "Cloud-Native Architecture",
        key_points=[
            "Separation of storage and compute for cost optimization",
            "Multi-cluster warehouses for workload isolation",
            "Cross-cloud and cross-region data replication"
        ]
    )

    composer.add_emphasis_slide(
        "Key Insight",
        quote="The future of data is cloud-native, elastic, and collaborative"
    )

    composer.save()
