"""
Create Sample PowerPoint Presentation
For testing the video automation system
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed")
    print("Install with: pip install python-pptx")
    exit(1)

def create_sample_presentation():
    """Create a simple test presentation"""

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]

    title1.text = "Welcome to Video Automation"
    subtitle1.text = "Testing the Chatterbox Video Creator"

    # Add speaker notes
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "Hello everyone! Welcome to this demonstration of our automated video creation system. I'm excited to show you how this works."

    # Slide 2: Features
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title2 = slide2.shapes.title
    title2.text = "Key Features"

    # Add content
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "Automatic voice narration"

    p = tf.add_paragraph()
    p.text = "PowerPoint slide rendering"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Audio-video synchronization"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "B-roll integration"
    p.level = 0

    # Add speaker notes
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "This system has four main features. First is automatic voice narration using AI. Second is professional slide rendering. Third is smart audio and video synchronization. And fourth is support for screen recordings and B-roll footage."

    # Slide 3: How It Works
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "How It Works"

    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Create PowerPoint with speaker notes"

    p = tf3.add_paragraph()
    p.text = "Run the video creator"
    p.level = 0

    p = tf3.add_paragraph()
    p.text = "Get YouTube-ready video"
    p.level = 0

    # Add speaker notes with pause
    notes_slide3 = slide3.notes_slide
    notes_slide3.notes_text_frame.text = "The process is simple. First, you create a PowerPoint presentation with your content and speaker notes. Then, you run our video creator tool. And finally, you get a complete YouTube-ready video with voice narration."

    # Slide 4: Demo
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Live Demo"

    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "This is a test demonstration"

    p = tf4.add_paragraph()
    p.text = "Generating video automatically"
    p.level = 0

    p = tf4.add_paragraph()
    p.text = "With AI voice narration"
    p.level = 0

    # Add speaker notes
    notes_slide4 = slide4.notes_slide
    notes_slide4.notes_text_frame.text = "Right now, you're watching this system in action. The video you're seeing was generated automatically from this PowerPoint presentation. Pretty cool, right?"

    # Slide 5: Thank You
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    title5 = slide5.shapes.title
    subtitle5 = slide5.placeholders[1]

    title5.text = "Thank You!"
    subtitle5.text = "Ready to create amazing videos"

    # Add speaker notes
    notes_slide5 = slide5.notes_slide
    notes_slide5.notes_text_frame.text = "Thank you for watching this demonstration! You're now ready to create your own automated videos. Happy creating!"

    # Save presentation
    output_file = "_video_automation/examples/sample_presentation.pptx"
    prs.save(output_file)

    print("Sample presentation created!")
    print(f"   File: {output_file}")
    print(f"   Slides: 5")
    print("\nYou can now generate a video with:")
    print(f"   python _video_automation/video_creator.py {output_file}")

    return output_file

if __name__ == "__main__":
    print("=" * 70)
    print("Creating Sample Presentation")
    print("=" * 70)
    print()

    create_sample_presentation()
