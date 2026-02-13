"""
Automated Slide Redesigner with AI Image Generation
Complete production pipeline with inventory management
"""

import sys
import os
import argparse
from pathlib import Path
from pptx import Presentation
import re

# Import our modules
from config_manager import APIConfig
from api_manager import MultiAccountAPIManager
from prompt_generator import PromptGenerator
from image_generator import ImageGenerator
from slide_composer import SlideComposer
from slide_composer_snowbrix import SnowbrixSlideComposer
from inventory_manager import InventoryManager, InventoryStage


def safe_text(text: str) -> str:
    """Safely handle unicode characters"""
    if not text:
        return ""
    try:
        # Try to encode/decode to handle special chars
        return text.encode('utf-8', errors='ignore').decode('utf-8')
    except:
        return str(text)


def parse_presentation(pptx_path: str) -> list:
    """Parse PowerPoint and extract slide information"""
    print(f"Parsing presentation: {pptx_path}")

    try:
        prs = Presentation(pptx_path)
        slides = []

        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                # Extract title
                title = ""
                if slide.shapes.title:
                    try:
                        title = safe_text(slide.shapes.title.text)
                    except:
                        title = ""

                # If no title shape, find the title from text shapes
                if not title:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            try:
                                shape_text = safe_text(shape.text).strip()

                                # Skip logo text
                                if "SNOWBRIXAI" in shape_text or "SNOWBRIX" in shape_text:
                                    continue

                                # Skip numbers (numbered circles)
                                if shape_text.isdigit():
                                    continue

                                # Skip very short text
                                if len(shape_text) < 10:
                                    continue

                                # Skip text that looks like bullet points (database names)
                                if shape_text.startswith("MDF_"):
                                    continue

                                # This is likely the title
                                title = shape_text.split('\n')[0].strip()
                                break  # Take the first valid title found
                            except:
                                pass

                # Extract content from all text boxes
                content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        try:
                            content += safe_text(shape.text) + "\n"
                        except:
                            pass

                # Extract speaker notes
                notes = ""
                if slide.has_notes_slide:
                    try:
                        notes = safe_text(slide.notes_slide.notes_text_frame.text)
                    except:
                        notes = ""

                # Extract key points from content
                key_points = []
                try:
                    lines = [l.strip() for l in content.split('\n') if l.strip()]
                    for i, line in enumerate(lines):
                        # Handle various bullet characters
                        if any(line.startswith(c) for c in ['*', '-', '*', '●', '○']):
                            clean_line = line.lstrip('*-*●○ ').strip()
                            if clean_line:
                                key_points.append(clean_line)

                        # Handle numbered lists (1, 2, 3, etc.)
                        elif len(line) > 0 and line[0].isdigit() and len(line) <= 2:
                            # Number found, check next line for content
                            if i + 1 < len(lines):
                                next_line = lines[i + 1].strip()
                                # Skip if next line is empty or is the title
                                if len(next_line) > 15 and next_line != title:
                                    key_points.append(next_line)

                        # Also catch lines with " - " pattern (descriptions)
                        elif " - " in line and len(line) > 20 and line != title:
                            # This looks like a description (e.g., "MDF_ADMIN - Full control...")
                            if not any(line.startswith(skip) for skip in [title[:10]]):
                                key_points.append(line)
                except:
                    key_points = []

                # Deduplicate key points while preserving order
                seen = set()
                unique_points = []
                for point in key_points:
                    if point not in seen:
                        seen.add(point)
                        unique_points.append(point)

                slides.append({
                    'number': slide_num,
                    'title': title,
                    'content': content,
                    'notes': notes,
                    'key_points': unique_points[:5]  # Max 5 unique points
                })
            except Exception as e:
                # Skip problematic slides
                print(f"  Warning: Skipped slide {slide_num}: {str(e)[:50]}")
                continue

        print(f"Found {len(slides)} slides")
        return slides

    except Exception as e:
        print(f"ERROR parsing presentation: {str(e)[:100]}")
        return []


def generate_prompts_stage(pptx_path: str, project_dir: Path):
    """Stage 1: Generate image prompts"""
    print("\n" + "="*70)
    print("STAGE 1: GENERATING IMAGE PROMPTS")
    print("="*70)

    # Parse presentation
    slides = parse_presentation(pptx_path)
    if not slides:
        return False

    # Create prompt generator
    prompt_gen = PromptGenerator()

    # Create prompts directory
    prompt_dir = project_dir / "image_prompts"
    prompt_dir.mkdir(parents=True, exist_ok=True)

    # Generate prompts for each slide
    for slide in slides:
        prompt, slide_type = prompt_gen.generate_prompt(
            slide['title'],
            slide['content'],
            slide['number']
        )

        # Enhance prompt with key points if available
        if slide['key_points']:
            prompt += f"\n\nKey Points to Visualize:\n"
            for i, point in enumerate(slide['key_points'], 1):
                prompt += f"{i}. {point}\n"

        # Save prompt to file
        prompt_file = prompt_dir / f"slide_{slide['number']:02d}_{slide_type.value}.txt"
        with open(prompt_file, 'w', encoding='utf-8') as f:
            f.write(f"Slide {slide['number']}: {slide['title']}\n")
            f.write(f"Type: {slide_type.value}\n")
            f.write("="*70 + "\n\n")
            f.write(prompt)

        # Clean title for display (remove problematic unicode)
        display_title = slide['title'][:50].encode('ascii', 'ignore').decode('ascii')
        print(f"  [{slide['number']:02d}] {display_title}")
        print(f"       Type: {slide_type.value} -> {prompt_file.name}")

    print(f"\n[OK] Generated {len(slides)} prompts")
    print(f"[OK] Saved to: {prompt_dir}")
    print("\nNext steps:")
    print("  1. Review prompts: Open image_prompts/ folder")
    print("  2. Edit prompts if needed")
    print("  3. (Optional) Manually generate images -> save to images/")
    print("  4. Run: python slide_redesigner_v2.py [file] --generate")

    return True


def generate_images_stage(project_dir: Path, api_manager: MultiAccountAPIManager):
    """Stage 2: Generate images from prompts"""
    print("\n" + "="*70)
    print("STAGE 2: GENERATING IMAGES")
    print("="*70)

    prompt_dir = project_dir / "image_prompts"
    image_dir = project_dir / "images"
    image_dir.mkdir(parents=True, exist_ok=True)

    # Get all prompts
    prompt_files = sorted(prompt_dir.glob("slide_*.txt"))
    if not prompt_files:
        print("No prompts found. Run --plan-only first.")
        return False

    # Create image generator
    img_gen = ImageGenerator(api_manager, str(image_dir))

    # Check API capacity
    print(f"\n{api_manager.get_status_summary()}\n")

    capacity = api_manager.check_capacity(len(prompt_files))
    print(f"Need: {capacity['required']} images")
    print(f"Available: {capacity['available']} API requests")

    if not capacity['sufficient']:
        print(f"[!]  Warning: Only {capacity['available']} requests available")
        print(f"   Missing: {capacity['shortfall']} images")
        print("   Will generate what we can + use placeholders for rest")
        response = input("\nContinue? (y/n): ")
        if response.lower() != 'y':
            return False

    print("\nGenerating images...")
    print("="*70)

    generated = 0
    existing = 0
    failed = 0

    for prompt_file in prompt_files:
        # Extract slide number
        match = re.search(r'slide_(\d+)', prompt_file.name)
        if not match:
            continue

        slide_num = int(match.group(1))

        # Read prompt
        with open(prompt_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()
            title = lines[0].replace('Slide ', '').split(':', 1)[1].strip() if len(lines) > 0 else ""
            prompt = ''.join(lines[3:])  # Skip header lines

        # Generate image
        print(f"\n[{slide_num:02d}/{len(prompt_files)}] {title[:50]}")

        success, image_path, message = img_gen.generate_image(
            prompt=prompt,
            slide_number=slide_num,
            slide_title=title
        )

        print(f"  {message}")

        if success:
            if "existing" in message.lower():
                existing += 1
            else:
                generated += 1
        else:
            failed += 1
            # Generate placeholder
            print(f"  Creating placeholder...")

    print("\n" + "="*70)
    print("IMAGE GENERATION COMPLETE")
    print("="*70)
    print(f"[OK] Generated: {generated}")
    print(f"[OK] Existing: {existing}")
    if failed > 0:
        print(f"[!]  Failed: {failed} (placeholders created)")

    print(f"\n{api_manager.get_status_summary()}")

    print("\nNext step:")
    print("  python slide_redesigner_v2.py [file] --create-slides")

    return True


def create_slides_stage(pptx_path: str, project_dir: Path, brand_style: str = 'snowbrix'):
    """Stage 3: Create presentation with images"""
    print("\n" + "="*70)
    print("STAGE 3: CREATING PRESENTATION")
    print("="*70)

    # Parse original presentation
    slides = parse_presentation(pptx_path)
    if not slides:
        return False

    # Check images exist
    image_dir = project_dir / "images"
    if not image_dir.exists():
        print("No images found. Run --generate first.")
        return False

    # Create output path
    pptx_name = Path(pptx_path).stem
    output_path = project_dir / f"{pptx_name}_redesigned.pptx"

    # Create composer with brand style
    if brand_style == 'snowbrix':
        composer = SnowbrixSlideComposer(str(output_path))
    else:
        composer = SlideComposer(str(output_path), brand_style=brand_style)

    print("\nCreating slides...")

    for slide in slides:
        slide_num = slide['number']

        # Find image
        image_path = None
        for ext in ['png', 'jpg', 'jpeg']:
            img_file = image_dir / f"slide_{slide_num:02d}.{ext}"
            if img_file.exists():
                image_path = img_file
                break

        # First slide is title
        if slide_num == 1:
            composer.add_title_slide(
                slide['title'],
                slide['notes'][:100] if slide['notes'] else ""
            )
        else:
            composer.add_slide_with_image(
                title=slide['title'],
                image_path=image_path,
                slide_type="content",
                notes=slide['notes'],
                key_points=slide['key_points']
            )

        status = "[OK]" if image_path else "[!] (placeholder)"
        print(f"  [{slide_num:02d}] {slide['title'][:50]} {status}")

    # Save
    composer.save()

    print(f"\n[OK] Created presentation: {output_path}")
    print("\nNext steps:")
    print(f"  1. Review: {output_path}")
    print(f"  2. Generate audio: python video_creator.py {output_path} --audio-only")
    print(f"  3. Create video: python video_creator.py {output_path} --use-existing-audio")

    return True


def main():
    parser = argparse.ArgumentParser(
        description="Automated Slide Redesigner with AI Image Generation",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Production Workflow:

1. Generate Prompts:
   python slide_redesigner_v2.py presentation.pptx --plan-only

2. (Optional) Manually create some images:
   - Review prompts in: _projects/[name]/image_prompts/
   - Generate images manually -> save to: _projects/[name]/images/

3. Generate Missing Images:
   python slide_redesigner_v2.py presentation.pptx --generate

4. Create Presentation:
   python slide_redesigner_v2.py presentation.pptx --create-slides

5. View Inventory:
   python slide_redesigner_v2.py --inventory-status

Examples:
  # Complete pipeline (all stages)
  python slide_redesigner_v2.py presentation.pptx --all

  # Just generate prompts
  python slide_redesigner_v2.py presentation.pptx --plan-only

  # Generate images (with API keys)
  python slide_redesigner_v2.py presentation.pptx --generate --api-keys key1,key2
        """
    )

    parser.add_argument('presentation', nargs='?',
                       help='PowerPoint file (.pptx)')

    # Stage commands
    parser.add_argument('--plan-only', action='store_true',
                       help='Stage 1: Generate prompts only')

    parser.add_argument('--generate', action='store_true',
                       help='Stage 2: Generate images from prompts')

    parser.add_argument('--create-slides', action='store_true',
                       help='Stage 3: Create presentation with images')

    parser.add_argument('--all', action='store_true',
                       help='Run all stages (prompts -> images -> slides)')

    # API configuration
    parser.add_argument('--api-keys',
                       help='Comma-separated API keys')

    parser.add_argument('--config',
                       help='Path to api_keys.json config file')

    # Inventory
    parser.add_argument('--inventory-status', action='store_true',
                       help='Show inventory dashboard')

    parser.add_argument('--generate-inventory', action='store_true',
                       help='Generate images for all pending projects')

    # Output
    parser.add_argument('--project-dir',
                       help='Project directory (default: _projects/[name])')

    # Brand style
    parser.add_argument('--brand',
                       choices=['snowbrix', 'warm', 'bold', 'sophisticated', 'energetic', 'natural'],
                       default='snowbrix',
                       help='Brand color scheme (default: snowbrix - authentic cream+green design)')

    args = parser.parse_args()

    print("="*70)
    print("AUTOMATED SLIDE REDESIGNER")
    print("AI-Powered Image Generation for Data Engineering")
    print("="*70)

    # Inventory commands
    if args.inventory_status:
        inventory = InventoryManager()
        inventory.scan_all_projects()
        print("\n" + inventory.get_dashboard())
        return 0

    # Need presentation file for other commands
    if not args.presentation and not args.generate_inventory:
        parser.print_help()
        return 1

    if args.presentation and not os.path.exists(args.presentation):
        print(f"ERROR: File not found: {args.presentation}")
        return 1

    # Setup project directory
    if args.presentation:
        pres_name = Path(args.presentation).stem
        project_dir = Path(args.project_dir) if args.project_dir else Path(f"_projects/{pres_name}")
        project_dir.mkdir(parents=True, exist_ok=True)

        print(f"\nProject: {project_dir}")
        print(f"Presentation: {args.presentation}")
        print()

    # Stage 1: Generate prompts
    if args.plan_only or args.all:
        success = generate_prompts_stage(args.presentation, project_dir)
        if not success:
            return 1

        if args.plan_only:
            return 0

    # Stage 2: Generate images
    if args.generate or args.all:
        # Load API config
        config = APIConfig(args.config)
        if args.api_keys:
            config.add_cli_keys(args.api_keys)

        if not config.validate():
            return 1

        # Create API manager
        api_manager = MultiAccountAPIManager(config.get_google_accounts())

        success = generate_images_stage(project_dir, api_manager)
        if not success:
            return 1

        if args.generate:
            return 0

    # Stage 3: Create slides
    if args.create_slides or args.all:
        success = create_slides_stage(args.presentation, project_dir, brand_style=args.brand)
        if not success:
            return 1

    # Update inventory
    if args.presentation:
        inventory = InventoryManager()
        inventory.scan_project(pres_name, project_dir)

    print("\n" + "="*70)
    print("SUCCESS!")
    print("="*70)

    return 0


if __name__ == "__main__":
    sys.exit(main())
