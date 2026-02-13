"""
GrowthSchool Slide Composer
Using authentic GrowthSchool color palette
"""

from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from brand_colors_growthschool import GrowthSchoolColors


class GrowthSchoolSlideComposer:
    """Slide composer with authentic GrowthSchool design"""

    def __init__(self, output_path: str):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # Get authentic GrowthSchool colors
        scheme = GrowthSchoolColors.get_scheme()
        self.colors = scheme
        print(f"Using: {scheme['name']}")

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Title slide - GrowthSchool style"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Large subtle green shape (top right)
        shape = slide.shapes.add_shape(
            1, Inches(12), Inches(-0.5),
            Inches(5), Inches(5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['color_primary_light']
        shape.fill.transparency = 0.5
        shape.line.fill.background()
        shape.rotation = 20

        # Thick Forest Green accent bar (left)
        bar = slide.shapes.add_shape(
            1, Inches(1.5), Inches(3.5),
            Inches(0.3), Inches(3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['color_primary']
        bar.line.fill.background()

        # Title in FOREST GREEN (#2D5F3F)
        title_box = slide.shapes.add_textbox(
            Inches(2.2), Inches(3.8),
            Inches(12), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']  # Forest Green

        # Subtitle in SAGE GREEN (#6B9080)
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(2.2), Inches(6),
                Inches(12), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.LEFT
            sub_para.font.size = Pt(36)
            sub_para.font.color.rgb = self.colors['color_secondary']  # Sage Green

        return slide

    def add_content_slide(
        self,
        title: str,
        key_points: List[str] = None,
        notes: str = ""
    ):
        """Content slide - GrowthSchool style with strong contrast"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Light cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_secondary']  # #FAF5F0

        # === TITLE SECTION ===
        # Subtle green background for title area
        title_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(16), Inches(2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['bg_accent']  # Light green #E8F3ED
        title_bg.line.fill.background()

        # Thick accent bar (Sage Green)
        bar = slide.shapes.add_shape(
            1, Inches(0.8), Inches(0.5),
            Inches(0.2), Inches(1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['color_secondary']  # Sage Green
        bar.line.fill.background()

        # Title in FOREST GREEN (#2D5F3F) - DARK & BOLD
        title_box = slide.shapes.add_textbox(
            Inches(1.3), Inches(0.6),
            Inches(14), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']  # Forest Green - STANDS OUT!

        # === KEY POINTS with numbered circles ===
        if key_points:
            start_y = 3
            for i, point in enumerate(key_points[:3]):
                y = start_y + (i * 1.6)

                # Circle with SAGE GREEN (#6B9080)
                circle = slide.shapes.add_shape(
                    1, Inches(1.5), Inches(y),
                    Inches(0.7), Inches(0.7)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']  # Sage Green circle
                circle.line.fill.background()

                # Number (cream text on green)
                num_box = slide.shapes.add_textbox(
                    Inches(1.5), Inches(y),
                    Inches(0.7), Inches(0.7)
                )
                num_frame = num_box.text_frame
                num_frame.clear()  # Clear any default formatting
                num_para = num_frame.paragraphs[0]
                num_para.text = str(i + 1)  # Set text via paragraph, not frame
                num_para.alignment = PP_ALIGN.CENTER
                num_para.font.size = Pt(36)
                num_para.font.bold = True
                num_para.font.color.rgb = self.colors['text_on_green']  # Cream on green

                # Point text in WARM GRAY (#4A4A4A) - READABLE & DARK
                point_box = slide.shapes.add_textbox(
                    Inches(2.5), Inches(y + 0.05),
                    Inches(12.5), Inches(0.9)
                )
                point_frame = point_box.text_frame
                point_frame.clear()  # Clear any default formatting
                point_frame.word_wrap = True
                point_para = point_frame.paragraphs[0]
                point_para.text = point[:120]  # Set text via paragraph, not frame
                point_para.font.size = Pt(30)
                point_para.font.bold = True
                point_para.font.color.rgb = self.colors['text_secondary']  # Warm Gray - DARK!

                # Explicitly remove bullet formatting at XML level
                from lxml.etree import Element
                pPr = point_para._element.get_or_add_pPr()
                # Remove any existing bullet elements
                for child in list(pPr):
                    if 'bu' in child.tag.lower():
                        pPr.remove(child)
                # Add buNone element to explicitly disable bullets
                buNone = Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                pPr.insert(0, buNone)

        # Add notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_emphasis_slide(self, title: str, quote: str = ""):
        """Emphasis/Quote slide with light green background"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Light green background (#E8F3ED)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_accent']

        # Large quote or emphasis text (centered, with margins)
        text_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5),
            Inches(13), Inches(4)
        )
        text_frame = text_box.text_frame
        text_frame.text = quote or title
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.alignment = PP_ALIGN.CENTER
        text_para.font.size = Pt(44)
        text_para.font.bold = True
        text_para.font.color.rgb = self.colors['color_primary']  # Forest Green

        return slide

    def add_slide_with_image(
        self,
        title: str,
        image_path=None,
        slide_type: str = "content",
        notes: str = "",
        key_points: list = None
    ):
        """Compatibility method for slide redesigner integration"""
        if slide_type == "emphasis" or slide_type == "quote":
            return self.add_emphasis_slide(title, notes)
        else:
            # Standard content slide
            return self.add_content_slide(title, key_points or [], notes)

    def save(self):
        """Save presentation"""
        self.prs.save(self.output_path)
        print(f"[OK] Saved: {self.output_path}")


# Demo
if __name__ == "__main__":
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    composer = GrowthSchoolSlideComposer(f"_AUTHENTIC_GrowthSchool_{timestamp}.pptx")

    composer.add_title_slide(
        "Data Engineering Mastery",
        "Transform Data into Strategic Advantage"
    )

    composer.add_content_slide(
        "Why Snowflake Wins",
        key_points=[
            "10x faster queries than traditional data warehouses",
            "Zero infrastructure management - fully automated scaling",
            "Instant data sharing across teams and organizations"
        ]
    )

    composer.add_content_slide(
        "Cloud-Native Architecture",
        key_points=[
            "Separation of storage and compute for cost optimization",
            "Multi-cluster warehouses for workload isolation",
            "Cross-cloud and cross-region data replication"
        ]
    )

    composer.add_emphasis_slide(
        "Key Insight",
        quote="The future of data is cloud-native, elastic, and collaborative"
    )

    composer.save()
