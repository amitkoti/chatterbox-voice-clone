"""
Enhanced Slide Composer with Strong Visual Impact
Makes colors POP on content slides
"""

from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

try:
    from brand_colors_refined import RefinedBrandColors
    REFINED_AVAILABLE = True
except:
    REFINED_AVAILABLE = False
from brand_colors import BrandColors


class SlideComposerV2:
    """Enhanced slide composer with strong visual contrast"""

    def __init__(self, output_path: str, brand_style: str = 'warm'):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # Get brand colors
        if REFINED_AVAILABLE and brand_style in ['warm', 'bold', 'sophisticated', 'energetic', 'natural']:
            brand_scheme = RefinedBrandColors.get_scheme(brand_style)
            self.colors = {
                'bg_primary': brand_scheme['bg_primary'],
                'bg_secondary': brand_scheme['bg_secondary'],
                'bg_card': brand_scheme['bg_card'],
                'text_primary': brand_scheme['text_primary'],
                'text_secondary': brand_scheme['text_secondary'],
                'text_accent': brand_scheme['text_accent'],
                'color_primary': brand_scheme['color_primary'],
                'color_accent': brand_scheme['color_accent'],
                'color_primary_dark': brand_scheme['color_primary_dark'],
                'color_primary_light': brand_scheme['color_primary_light'],
            }
        else:
            brand_scheme = BrandColors.get_scheme(brand_style)
            self.colors = brand_scheme

        self.brand_name = brand_scheme['name']
        print(f"Using brand style: {self.brand_name}")

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Title slide with strong brand presence"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Large colored accent shape (top right corner)
        accent_shape = slide.shapes.add_shape(
            1, Inches(11), Inches(-1),
            Inches(6), Inches(6)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = self.colors['color_accent']
        accent_shape.fill.transparency = 0.15
        accent_shape.line.fill.background()
        accent_shape.rotation = 25

        # Thick vertical accent bar (left)
        bar = slide.shapes.add_shape(
            1, Inches(1), Inches(3),
            Inches(0.25), Inches(3.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['color_primary']
        bar.line.fill.background()

        # Title in PRIMARY color (dark, stands out)
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.5),
            Inches(13), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.size = Pt(76)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Subtitle in ACCENT color (bright, pops)
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(6),
                Inches(12), Inches(1)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.LEFT
            sub_para.font.size = Pt(36)
            sub_para.font.color.rgb = self.colors['color_accent']

        return slide

    def add_content_slide(
        self,
        title: str,
        key_points: List[str] = None,
        notes: str = ""
    ):
        """Content slide with strong visual hierarchy"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Cream background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # TITLE SECTION with colored background
        # Add colored rectangle behind title for impact
        title_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(16), Inches(2.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['color_primary']
        title_bg.fill.transparency = 0.05
        title_bg.line.fill.background()

        # Thick accent bar on left of title
        title_bar = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0.5),
            Inches(0.2), Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.colors['color_accent']
        title_bar.line.fill.background()

        # Title in dark PRIMARY color
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.6),
            Inches(14), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # KEY POINTS with colored accent cards
        if key_points:
            start_y = 3
            for i, point in enumerate(key_points[:3]):
                y = start_y + (i * 1.5)

                # Colored number circle
                num_circle = slide.shapes.add_shape(
                    1, Inches(1.2), Inches(y + 0.1),
                    Inches(0.6), Inches(0.6)
                )
                num_circle.fill.solid()
                num_circle.fill.fore_color.rgb = self.colors['color_accent']
                num_circle.line.fill.background()

                # Number text (white on colored circle)
                num_box = slide.shapes.add_textbox(
                    Inches(1.2), Inches(y + 0.1),
                    Inches(0.6), Inches(0.6)
                )
                num_frame = num_box.text_frame
                num_frame.text = str(i + 1)
                num_para = num_frame.paragraphs[0]
                num_para.alignment = PP_ALIGN.CENTER
                num_para.font.size = Pt(32)
                num_para.font.bold = True
                num_para.font.color.rgb = RGBColor(255, 255, 255)

                # Point text in PRIMARY color (dark, strong)
                point_box = slide.shapes.add_textbox(
                    Inches(2.2), Inches(y + 0.15),
                    Inches(12.5), Inches(0.8)
                )
                point_frame = point_box.text_frame
                point_frame.text = point[:120]
                point_para = point_frame.paragraphs[0]
                point_para.font.size = Pt(32)
                point_para.font.bold = True
                point_para.font.color.rgb = self.colors['color_primary']

        # Add notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def save(self):
        """Save presentation"""
        self.prs.save(self.output_path)
        print(f"[OK] Saved: {self.output_path}")


# Test/Demo
if __name__ == "__main__":
    composer = SlideComposerV2("_test_v2_sample.pptx", "warm")

    composer.add_title_slide(
        "Data Engineering Excellence",
        "Modern Cloud Solutions"
    )

    composer.add_content_slide(
        "Why Snowflake?",
        key_points=[
            "Lightning-fast query performance at any scale",
            "Zero maintenance - fully managed cloud service",
            "Secure data sharing across your organization"
        ]
    )

    composer.add_content_slide(
        "Key Benefits",
        key_points=[
            "Instant elasticity scales with your workload",
            "Time travel and data cloning capabilities",
            "Multi-cloud support (AWS, Azure, GCP)"
        ]
    )

    composer.save()
