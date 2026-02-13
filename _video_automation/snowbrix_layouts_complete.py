"""
Complete Snowbrix Layout System
All professional layouts from Snowflake template adapted for Snowbrix
"""

from pathlib import Path
from typing import List, Optional, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from brand_colors_snowbrix import SnowbrixColors
from lxml.etree import Element


class SnowbrixLayoutsComplete:
    """Complete professional layout system for Snowbrix"""

    def __init__(self, output_path: str, include_logo: bool = True, include_page_numbers: bool = False):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.include_logo = include_logo
        self.include_page_numbers = include_page_numbers
        self.slide_count = 0

        # Get Snowbrix colors
        scheme = SnowbrixColors.get_scheme()
        self.colors = scheme
        print(f"Using: {scheme['name']}")

    def _remove_bullets(self, paragraph):
        """Remove bullet formatting from paragraph"""
        try:
            pPr = paragraph._element.get_or_add_pPr()
            for child in list(pPr):
                if 'bu' in child.tag.lower():
                    pPr.remove(child)
            buNone = Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
            pPr.insert(0, buNone)
        except:
            pass

    def _add_logo_to_slide(self, slide):
        """Add Snowbrixai logo to bottom right"""
        if not self.include_logo:
            return

        logo_x, logo_y = 13.7, 8.3
        font_size = Pt(22)

        logo_box = slide.shapes.add_textbox(
            Inches(logo_x), Inches(logo_y),
            Inches(2.5), Inches(0.5)
        )

        text_frame = logo_box.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]

        # SNOW (forest green)
        run_snow = paragraph.add_run()
        run_snow.text = "SNOW"
        run_snow.font.size = font_size
        run_snow.font.bold = True
        run_snow.font.color.rgb = self.colors['color_primary']

        # BRIX (terracotta)
        run_brix = paragraph.add_run()
        run_brix.text = "BRIX"
        run_brix.font.size = font_size
        run_brix.font.bold = True
        run_brix.font.color.rgb = self.colors['accent_terracotta']

        # AI (charcoal)
        run_ai = paragraph.add_run()
        run_ai.text = "AI"
        run_ai.font.size = font_size
        run_ai.font.bold = True
        run_ai.font.color.rgb = self.colors['text_secondary']

    def _add_page_number(self, slide):
        """Add page number to bottom right"""
        if not self.include_page_numbers:
            return

        self.slide_count += 1

        page_box = slide.shapes.add_textbox(
            Inches(15.2), Inches(8.5),
            Inches(0.6), Inches(0.3)
        )
        page_frame = page_box.text_frame
        page_frame.text = str(self.slide_count)
        page_para = page_frame.paragraphs[0]
        page_para.alignment = PP_ALIGN.RIGHT
        page_para.font.size = Pt(18)
        page_para.font.color.rgb = self.colors['color_secondary']

    # ========================================================================
    # LAYOUT 1: TITLE SLIDE
    # ========================================================================

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Title slide - Hero introduction"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Decorative green shape (top right - stays within bounds)
        shape = slide.shapes.add_shape(
            1, Inches(13), Inches(0.3),
            Inches(2.8), Inches(2.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['color_primary_light']
        shape.fill.transparency = 0.6
        shape.line.fill.background()
        # No rotation to keep within boundaries

        # Accent bar (left)
        bar = slide.shapes.add_shape(
            1, Inches(1.5), Inches(3.5),
            Inches(0.3), Inches(3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['color_primary']
        bar.line.fill.background()

        # Title (with word wrapping)
        title_box = slide.shapes.add_textbox(
            Inches(2.2), Inches(3.8),
            Inches(11.5), Inches(2.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.size = Pt(68)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(2.2), Inches(6),
                Inches(12), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.LEFT
            sub_para.font.size = Pt(36)
            sub_para.font.color.rgb = self.colors['color_secondary']

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    # ========================================================================
    # LAYOUT 2: CONTENT WITH BULLETS (3-5 points)
    # ========================================================================

    def add_content_slide(self, title: str, key_points: List[str] = None,
                         subtitle: str = "", notes: str = ""):
        """Content slide with numbered bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_secondary']

        # Title section
        title_bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(16), Inches(2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['bg_accent']
        title_bg.line.fill.background()

        # Accent bar
        bar = slide.shapes.add_shape(
            1, Inches(0.8), Inches(0.5),
            Inches(0.2), Inches(1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['color_secondary']
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.3), Inches(0.6),
            Inches(14), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Subtitle (if provided)
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(1.5),
                Inches(14), Inches(0.5)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = self.colors['color_secondary']

        # Key points
        if key_points:
            num_points = min(len(key_points), 5)
            if num_points <= 3:
                start_y, spacing = 3.0, 1.6
            elif num_points == 4:
                start_y, spacing = 2.8, 1.3
            else:
                start_y, spacing = 2.6, 1.15

            for i, point in enumerate(key_points[:5]):
                y = start_y + (i * spacing)

                # Circle
                circle = slide.shapes.add_shape(
                    1, Inches(1.5), Inches(y),
                    Inches(0.7), Inches(0.7)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Number
                num_box = slide.shapes.add_textbox(
                    Inches(1.5), Inches(y),
                    Inches(0.7), Inches(0.7)
                )
                num_frame = num_box.text_frame
                num_frame.clear()
                num_para = num_frame.paragraphs[0]
                num_para.text = str(i + 1)
                num_para.alignment = PP_ALIGN.CENTER
                num_para.font.size = Pt(36)
                num_para.font.bold = True
                num_para.font.color.rgb = self.colors['text_on_green']

                # Point text
                point_box = slide.shapes.add_textbox(
                    Inches(2.5), Inches(y + 0.05),
                    Inches(12.5), Inches(0.9)
                )
                point_frame = point_box.text_frame
                point_frame.clear()
                point_frame.word_wrap = True
                point_para = point_frame.paragraphs[0]
                point_para.text = point[:120]
                point_para.font.size = Pt(30)
                point_para.font.bold = True
                point_para.font.color.rgb = self.colors['text_secondary']
                self._remove_bullets(point_para)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    # ========================================================================
    # LAYOUT 3: THREE-COLUMN LAYOUT
    # ========================================================================

    def add_three_column_slide(self, title: str,
                               col1_title: str = "", col1_points: List[str] = None,
                               col2_title: str = "", col2_points: List[str] = None,
                               col3_title: str = "", col3_points: List[str] = None,
                               notes: str = ""):
        """Three-column layout for feature grids"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_secondary']

        # Title section
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(14), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Three columns
        columns = [
            (0.8, col1_title, col1_points),
            (5.8, col2_title, col2_points),
            (10.8, col3_title, col3_points)
        ]

        for col_x, col_title, col_points in columns:
            col_y = 1.8

            # Column title
            if col_title:
                ct_box = slide.shapes.add_textbox(
                    Inches(col_x), Inches(col_y),
                    Inches(4.5), Inches(0.5)
                )
                ct_frame = ct_box.text_frame
                ct_frame.text = col_title
                ct_para = ct_frame.paragraphs[0]
                ct_para.font.size = Pt(28)
                ct_para.font.bold = True
                ct_para.font.color.rgb = self.colors['color_secondary']
                col_y += 1.0  # Increased from 0.7 to 1.0

            # Column points with sage green circles
            if col_points:
                for i, point in enumerate(col_points[:4]):
                    item_y = col_y + i * 0.7  # Reduced from 1.0 to 0.7

                    # Green bullet circle
                    circle = slide.shapes.add_shape(
                        1, Inches(col_x), Inches(item_y + 0.05),
                        Inches(0.18), Inches(0.18)
                    )
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = self.colors['color_secondary']
                    circle.line.fill.background()

                    # Point text
                    point_box = slide.shapes.add_textbox(
                        Inches(col_x + 0.35), Inches(item_y),
                        Inches(4.15), Inches(0.6)
                    )
                    pf = point_box.text_frame
                    pf.clear()
                    pf.text = point
                    pf.word_wrap = True
                    pp = pf.paragraphs[0]
                    pp.font.size = Pt(20)
                    pp.font.color.rgb = self.colors['text_secondary']
                    self._remove_bullets(pp)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    # ========================================================================
    # LAYOUT 4: IMAGE + TEXT SPLIT
    # ========================================================================

    def add_image_text_split(self, title: str, image_path: str = None,
                            content_points: List[str] = None,
                            image_on_left: bool = True, notes: str = ""):
        """Image on one side, text on other - very common layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(14), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Calculate positions
        if image_on_left:
            image_x, text_x = 0.8, 8.5
        else:
            text_x, image_x = 0.8, 8.5

        image_y, image_w, image_h = 1.8, 7.0, 6.0
        text_w = 6.5

        # Add image or placeholder
        if image_path and Path(image_path).exists():
            try:
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(image_x), Inches(image_y),
                    width=Inches(image_w), height=Inches(image_h)
                )
            except:
                self._add_image_placeholder(slide, image_x, image_y, image_w, image_h)
        else:
            self._add_image_placeholder(slide, image_x, image_y, image_w, image_h)

        # Add content points
        if content_points:
            for i, point in enumerate(content_points[:5]):
                y = 2.2 + (i * 1.1)

                # Bullet circle
                circle = slide.shapes.add_shape(
                    1, Inches(text_x), Inches(y),
                    Inches(0.4), Inches(0.4)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Point text
                point_box = slide.shapes.add_textbox(
                    Inches(text_x + 0.6), Inches(y - 0.05),
                    Inches(text_w - 0.7), Inches(0.9)
                )
                pf = point_box.text_frame
                pf.clear()
                pf.word_wrap = True
                pp = pf.paragraphs[0]
                pp.text = point
                pp.font.size = Pt(24)
                pp.font.color.rgb = self.colors['text_secondary']
                self._remove_bullets(pp)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    def _add_image_placeholder(self, slide, x, y, w, h):
        """Add image placeholder box"""
        placeholder = slide.shapes.add_shape(
            1, Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = self.colors['bg_accent']
        placeholder.line.color.rgb = self.colors['color_secondary']
        placeholder.line.width = Pt(2)

        # Placeholder text
        pt_box = slide.shapes.add_textbox(
            Inches(x + w/2 - 1.5), Inches(y + h/2 - 0.3),
            Inches(3), Inches(0.6)
        )
        ptf = pt_box.text_frame
        ptf.text = "[Image Placeholder]"
        ptp = ptf.paragraphs[0]
        ptp.alignment = PP_ALIGN.CENTER
        ptp.font.size = Pt(24)
        ptp.font.color.rgb = self.colors['color_secondary']

    # ========================================================================
    # LAYOUT 5: PARAGRAPH LAYOUT (Text-heavy)
    # ========================================================================

    def add_paragraph_slide(self, title: str, paragraphs: List[str],
                           subtitle: str = "", notes: str = ""):
        """Text-heavy layout with paragraphs instead of bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(14), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Subtitle
        start_y = 1.2
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(start_y),
                Inches(14), Inches(0.5)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.font.size = Pt(28)
            sub_para.font.color.rgb = self.colors['color_secondary']
            start_y += 0.7

        # Paragraphs
        if paragraphs:
            para_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(start_y),
                Inches(14), Inches(7 - start_y)
            )
            para_frame = para_box.text_frame
            para_frame.word_wrap = True
            para_frame.clear()

            for i, para_text in enumerate(paragraphs[:3]):
                if i > 0:
                    para_frame.add_paragraph()
                p = para_frame.paragraphs[i]
                p.text = para_text
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors['text_secondary']
                p.space_after = Pt(12)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    # ========================================================================
    # LAYOUT 6: AGENDA LAYOUT
    # ========================================================================

    def add_agenda_slide(self, title: str, agenda_items: List[Dict[str, str]], notes: str = ""):
        """Agenda layout with numbered items

        agenda_items format: [
            {"number": "01", "title": "Introduction", "duration": "5 min"},
            {"number": "02", "title": "Main Content", "duration": "15 min"},
            ...
        ]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5),
            Inches(14), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(52)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        # Agenda items (using numbered circles like content slides)
        if agenda_items:
            start_y = 2.0
            for i, item in enumerate(agenda_items[:6]):
                y = start_y + (i * 1.1)

                # Sage green circle (same style as content slides)
                circle = slide.shapes.add_shape(
                    1, Inches(1.0), Inches(y),
                    Inches(0.6), Inches(0.6)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Number INSIDE circle (white text, centered)
                num_text = item.get('number', str(i+1).zfill(2))
                num_box = slide.shapes.add_textbox(
                    Inches(1.0), Inches(y),
                    Inches(0.6), Inches(0.6)
                )
                nf = num_box.text_frame
                nf.clear()
                np = nf.paragraphs[0]
                np.text = num_text
                np.alignment = PP_ALIGN.CENTER
                np.font.size = Pt(28)
                np.font.bold = True
                np.font.color.rgb = self.colors['text_on_green']  # White text

                # Item title (aligned with circle)
                item_title = item.get('title', '')
                title_box = slide.shapes.add_textbox(
                    Inches(1.9), Inches(y + 0.05),
                    Inches(10), Inches(0.6)
                )
                tf = title_box.text_frame
                tf.text = item_title
                tf.word_wrap = True
                tp = tf.paragraphs[0]
                tp.font.size = Pt(32)
                tp.font.bold = True
                tp.font.color.rgb = self.colors['text_secondary']

                # Duration (right-aligned)
                duration = item.get('duration', '')
                if duration:
                    dur_box = slide.shapes.add_textbox(
                        Inches(12.5), Inches(y + 0.1),
                        Inches(2.5), Inches(0.5)
                    )
                    df = dur_box.text_frame
                    df.text = duration
                    dp = df.paragraphs[0]
                    dp.alignment = PP_ALIGN.RIGHT
                    dp.font.size = Pt(24)
                    dp.font.color.rgb = self.colors['color_secondary']

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    # ========================================================================
    # LAYOUT 7: TWO-COLUMN TEXT
    # ========================================================================

    def add_two_column_slide(self, title: str,
                            left_content: List[str] = None,
                            right_content: List[str] = None,
                            left_title: str = "",
                            right_title: str = "",
                            notes: str = ""):
        """Two-column text layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_secondary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(14), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        col_y = 1.8

        # Left column title
        if left_title:
            lt_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(col_y),  # Heading at left edge
                Inches(6.5), Inches(0.5)
            )
            ltf = lt_box.text_frame
            ltf.text = left_title
            ltp = ltf.paragraphs[0]
            ltp.font.size = Pt(32)
            ltp.font.bold = True
            ltp.font.color.rgb = self.colors['color_secondary']

        # Left content (bullets indented under heading)
        if left_content:
            content_y = col_y + (1.0 if left_title else 0)  # Increased from 0.7 to 1.0
            for i, item in enumerate(left_content[:5]):
                y_pos = content_y + i * 0.75  # Reduced from 1.1 to 0.75

                # Green bullet circle (indented under heading)
                circle = slide.shapes.add_shape(
                    1, Inches(1.0), Inches(y_pos + 0.1),  # Indented from 0.8 to 1.0
                    Inches(0.2), Inches(0.2)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Item text (indented further right)
                item_box = slide.shapes.add_textbox(
                    Inches(1.4), Inches(y_pos),  # Indented from 1.2 to 1.4
                    Inches(6.0), Inches(0.9)
                )
                itf = item_box.text_frame
                itf.clear()
                itf.text = item
                itf.word_wrap = True
                itp = itf.paragraphs[0]
                itp.font.size = Pt(24)
                itp.font.color.rgb = self.colors['text_secondary']
                self._remove_bullets(itp)

        # Right column title
        if right_title:
            rt_box = slide.shapes.add_textbox(
                Inches(8.5), Inches(col_y),  # Heading at column edge
                Inches(6.5), Inches(0.5)
            )
            rtf = rt_box.text_frame
            rtf.text = right_title
            rtp = rtf.paragraphs[0]
            rtp.font.size = Pt(32)
            rtp.font.bold = True
            rtp.font.color.rgb = self.colors['color_secondary']

        # Right content (bullets indented under heading)
        if right_content:
            content_y = col_y + (1.0 if right_title else 0)  # Increased from 0.7 to 1.0
            for i, item in enumerate(right_content[:5]):
                y_pos = content_y + i * 0.75  # Reduced from 1.1 to 0.75

                # Green bullet circle (indented under heading)
                circle = slide.shapes.add_shape(
                    1, Inches(8.7), Inches(y_pos + 0.1),  # Indented from 8.5 to 8.7
                    Inches(0.2), Inches(0.2)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.colors['color_secondary']
                circle.line.fill.background()

                # Item text (indented further right)
                item_box = slide.shapes.add_textbox(
                    Inches(9.1), Inches(y_pos),  # Indented from 8.9 to 9.1
                    Inches(6.0), Inches(0.9)
                )
                itf = item_box.text_frame
                itf.clear()
                itf.text = item
                itf.word_wrap = True
                itp = itf.paragraphs[0]
                itp.font.size = Pt(24)
                itp.font.color.rgb = self.colors['text_secondary']
                self._remove_bullets(itp)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    # ========================================================================
    # LAYOUT 8: TWO-COLUMN WITH PARAGRAPHS
    # ========================================================================

    def add_two_column_paragraph_slide(self, title: str,
                                      left_title: str = "", left_paragraphs: List[str] = None,
                                      right_title: str = "", right_paragraphs: List[str] = None,
                                      notes: str = ""):
        """Two-column with paragraph text instead of bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4),
            Inches(14), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['color_primary']

        col_y = 1.5

        # Left column
        if left_title:
            lt_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(col_y),
                Inches(7), Inches(0.5)
            )
            ltf = lt_box.text_frame
            ltf.text = left_title
            ltp = ltf.paragraphs[0]
            ltp.font.size = Pt(28)
            ltp.font.bold = True
            ltp.font.color.rgb = self.colors['color_secondary']

        if left_paragraphs:
            para_y = col_y + (0.7 if left_title else 0)
            left_para_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(para_y),
                Inches(7), Inches(6.5 - para_y)
            )
            lpf = left_para_box.text_frame
            lpf.word_wrap = True
            lpf.clear()

            for i, para_text in enumerate(left_paragraphs[:2]):
                if i > 0:
                    lpf.add_paragraph()
                p = lpf.paragraphs[i]
                p.text = para_text
                p.font.size = Pt(18)
                p.font.color.rgb = self.colors['text_secondary']
                p.space_after = Pt(12)

        # Right column
        if right_title:
            rt_box = slide.shapes.add_textbox(
                Inches(8.2), Inches(col_y),
                Inches(7), Inches(0.5)
            )
            rtf = rt_box.text_frame
            rtf.text = right_title
            rtp = rtf.paragraphs[0]
            rtp.font.size = Pt(28)
            rtp.font.bold = True
            rtp.font.color.rgb = self.colors['color_secondary']

        if right_paragraphs:
            para_y = col_y + (0.7 if right_title else 0)
            right_para_box = slide.shapes.add_textbox(
                Inches(8.2), Inches(para_y),
                Inches(7), Inches(6.5 - para_y)
            )
            rpf = right_para_box.text_frame
            rpf.word_wrap = True
            rpf.clear()

            for i, para_text in enumerate(right_paragraphs[:2]):
                if i > 0:
                    rpf.add_paragraph()
                p = rpf.paragraphs[i]
                p.text = para_text
                p.font.size = Pt(18)
                p.font.color.rgb = self.colors['text_secondary']
                p.space_after = Pt(12)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    # ========================================================================
    # EXISTING LAYOUTS (from previous implementation)
    # ========================================================================

    def add_emphasis_slide(self, title: str, quote: str = ""):
        """Emphasis/quote slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_accent']

        text_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5),
            Inches(13), Inches(4)
        )
        text_frame = text_box.text_frame
        text_frame.text = quote or title
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.alignment = PP_ALIGN.CENTER
        text_para.font.size = Pt(44)
        text_para.font.bold = True
        text_para.font.color.rgb = self.colors['color_primary']

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    def add_section_divider(self, section_title: str, subtitle: str = ""):
        """Section divider with forest green background"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['color_primary']

        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5),
            Inches(12), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(64)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_on_green']

        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(3), Inches(5.2),
                Inches(10), Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_para = sub_frame.paragraphs[0]
            sub_para.alignment = PP_ALIGN.CENTER
            sub_para.font.size = Pt(28)
            sub_para.font.color.rgb = self.colors['bg_primary']

        self._add_page_number(slide)
        return slide

    def add_thank_you_slide(self, message: str = "Thank You", contact_info: str = ""):
        """Thank you closing slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_accent']

        thanks_box = slide.shapes.add_textbox(
            Inches(2), Inches(3),
            Inches(12), Inches(2)
        )
        thanks_frame = thanks_box.text_frame
        thanks_frame.text = message
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.alignment = PP_ALIGN.CENTER
        thanks_para.font.size = Pt(72)
        thanks_para.font.bold = True
        thanks_para.font.color.rgb = self.colors['color_primary']

        if contact_info:
            contact_box = slide.shapes.add_textbox(
                Inches(3), Inches(5.5),
                Inches(10), Inches(1)
            )
            contact_frame = contact_box.text_frame
            contact_frame.text = contact_info
            contact_para = contact_frame.paragraphs[0]
            contact_para.alignment = PP_ALIGN.CENTER
            contact_para.font.size = Pt(24)
            contact_para.font.color.rgb = self.colors['color_secondary']

        self._add_logo_to_slide(slide)
        self._add_page_number(slide)
        return slide

    def save(self):
        """Save presentation"""
        self.prs.save(self.output_path)
        print(f"[OK] Saved: {self.output_path}")
