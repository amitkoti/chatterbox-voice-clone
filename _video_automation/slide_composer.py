"""
Professional Slide Composer
Creates visually stunning slides with images and minimal text
"""

from pathlib import Path
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
try:
    from brand_colors_refined import RefinedBrandColors
    REFINED_AVAILABLE = True
except:
    REFINED_AVAILABLE = False
from brand_colors import BrandColors


class SlideComposer:
    """Composes professional presentation slides"""

    def __init__(self, output_path: str, brand_style: str = 'warm'):
        self.output_path = output_path
        self.prs = Presentation()

        # Set 16:9 aspect ratio
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # Get brand color scheme (try refined first, fallback to original)
        if REFINED_AVAILABLE and brand_style in ['warm', 'bold', 'sophisticated', 'energetic', 'natural']:
            brand_scheme = RefinedBrandColors.get_scheme(brand_style)
            # Map refined colors to expected keys
            self.colors = {
                'bg_primary': brand_scheme['bg_primary'],
                'bg_secondary': brand_scheme['bg_secondary'],
                'bg_card': brand_scheme['bg_card'],
                'text_primary': brand_scheme['text_primary'],
                'text_secondary': brand_scheme['text_secondary'],
                'text_accent': brand_scheme['text_accent'],
                'green_primary': brand_scheme['color_primary'],
                'green_accent': brand_scheme['color_accent'],
                'green_dark': brand_scheme['color_primary_dark'],
                'green_light': brand_scheme['color_primary_light'],
            }
        else:
            brand_scheme = BrandColors.get_scheme(brand_style)
            self.colors = brand_scheme

        self.brand_name = brand_scheme['name']
        print(f"Using brand style: {self.brand_name}")

    def add_slide_with_image(
        self,
        title: str,
        image_path: Optional[Path],
        slide_type: str,
        notes: str = "",
        key_points: List[str] = None
    ):
        """Add a slide with image and minimal text"""

        # Use blank layout
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Brand background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['bg_primary']

        # Add image (full bleed or large)
        if image_path and image_path.exists():
            self._add_background_image(slide, image_path)

        # Add title overlay
        self._add_title_overlay(slide, title)

        # Add key points if provided (max 3)
        if key_points and len(key_points) > 0:
            self._add_key_points(slide, key_points[:3])

        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def _add_background_image(self, slide, image_path: Path):
        """Add image as background or large focal point"""
        try:
            # Get image dimensions
            img = PILImage.open(image_path)
            img_width, img_height = img.size
            img_ratio = img_width / img_height

            # Slide dimensions
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height
            slide_ratio = slide_width / slide_height

            # Calculate dimensions to fill slide
            if img_ratio > slide_ratio:
                # Image is wider - fit to height
                height = slide_height
                width = int(height * img_ratio)
                left = -(width - slide_width) / 2
                top = 0
            else:
                # Image is taller - fit to width
                width = slide_width
                height = int(width / img_ratio)
                left = 0
                top = -(height - slide_height) / 2

            # Add image
            slide.shapes.add_picture(
                str(image_path),
                left, top,
                width, height
            )

            # Add dark overlay for text readability
            overlay = slide.shapes.add_shape(
                1,  # Rectangle
                0, 0,
                slide_width, slide_height
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.3  # 30% transparent
            overlay.line.fill.background()

        except Exception as e:
            print(f"   Warning: Could not add image: {e}")

    def _add_title_overlay(self, slide, title: str):
        """Add title as overlay on image"""
        # Title box at top
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(15), Inches(1.5)
        )

        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']

        # Add subtle background to title for readability
        title_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.3), Inches(0.3),
            Inches(15.4), Inches(1.9)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
        title_bg.fill.transparency = 0.5
        title_bg.line.fill.background()

        # Move title box to front
        slide.shapes._spTree.remove(title_box._element)
        slide.shapes._spTree.append(title_box._element)

    def _add_key_points(self, slide, key_points: List[str]):
        """Add key bullet points at bottom"""
        points_text = "\n".join([f"• {point[:80]}" for point in key_points])

        points_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(14), Inches(2)
        )

        points_frame = points_box.text_frame
        points_frame.text = points_text
        points_frame.word_wrap = True

        for para in points_frame.paragraphs:
            para.font.size = Pt(24)
            para.font.color.rgb = self.colors['text_primary']
            para.space_after = Pt(8)

        # Background for points
        points_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(6.3),
            Inches(14.4), Inches(2.4)
        )
        points_bg.fill.solid()
        points_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
        points_bg.fill.transparency = 0.6
        points_bg.line.fill.background()

        # Move points to front
        slide.shapes._spTree.remove(points_box._element)
        slide.shapes._spTree.append(points_box._element)

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add title slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Brand background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['bg_primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(14), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(64)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5),
                Inches(14), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle

            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = self.colors['text_secondary']

        # Accent line
        line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(6), Inches(5.2),
            Inches(4), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['green_accent']
        line.line.fill.background()

        return slide

    def save(self):
        """Save presentation"""
        self.prs.save(self.output_path)
        print(f"[OK] Saved presentation: {self.output_path}")

    def get_slide_count(self) -> int:
        """Get number of slides"""
        return len(self.prs.slides)


def main():
    """Test slide composer"""
    composer = SlideComposer("test_output.pptx")

    # Add title slide
    composer.add_title_slide(
        "Data Engineering with Snowflake",
        "Modern Data Platform Architecture"
    )

    # Add content slide
    composer.add_slide_with_image(
        title="Pipeline Architecture",
        image_path=None,  # No image for test
        slide_type="architecture",
        notes="This is the speaker notes for this slide",
        key_points=[
            "Scalable data ingestion",
            "Real-time processing",
            "Cloud-native architecture"
        ]
    )

    composer.save()
    print(f"Created {composer.get_slide_count()} slides")


if __name__ == "__main__":
    main()
