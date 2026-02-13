"""
Update Module 01 presentation by adding full-screen zoom slides after key slides.

Keeps original slides intact and adds full-screen image slides for detailed viewing.
"""
import sys
import os
from pathlib import Path
from datetime import datetime

# Add video automation to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "_video_automation"))

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image

def remove_watermark_from_image(image_path, crop_percentage=10):
    """
    Crop the bottom portion of image to remove Gemini watermark.

    Args:
        image_path: Path to original image
        crop_percentage: Percentage of height to crop from bottom (default 10%)

    Returns:
        Path to cropped image
    """
    # Open image
    img = Image.open(image_path)
    width, height = img.size

    # Crop bottom percentage (where watermarks appear)
    crop_height = int(height * (100 - crop_percentage) / 100)

    cropped_img = img.crop((0, 0, width, crop_height))

    # Save as new file
    cropped_path = str(image_path).replace('.png', '_nowatermark.png')
    cropped_img.save(cropped_path, 'PNG', quality=95)

    print(f"       Cropped {crop_percentage}% from bottom: {cropped_path}")
    return cropped_path

def add_fullscreen_image_slide(prs, image_path, slide_position):
    """
    Insert a full-screen image slide at the specified position.

    Args:
        prs: Presentation object
        image_path: Path to image file
        slide_position: Index where to insert the slide

    Returns:
        The created slide object
    """
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Move slide to correct position
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])  # Remove from end
    xml_slides.insert(slide_position, slides_list[-1])  # Insert at position

    # Set black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Get image dimensions
    img = Image.open(image_path)
    img_width, img_height = img.size
    img_aspect = img_width / img_height

    # PowerPoint slide dimensions
    slide_width = 16.0
    slide_height = 9.0
    slide_aspect = slide_width / slide_height

    # Calculate positioning to fit image perfectly (maintain aspect ratio)
    if img_aspect > slide_aspect:
        # Image is wider - fit to width
        width = Inches(slide_width)
        height = Inches(slide_width / img_aspect)
        left = Inches(0)
        top = Inches((slide_height - (slide_width / img_aspect)) / 2)
    else:
        # Image is taller - fit to height
        height = Inches(slide_height)
        width = Inches(slide_height * img_aspect)
        left = Inches((slide_width - (slide_height * img_aspect)) / 2)
        top = Inches(0)

    # Add full-screen image
    picture = slide.shapes.add_picture(
        str(image_path),
        left, top,
        width=width,
        height=height
    )

    return slide

def main():
    print("=" * 70)
    print("Module 01 - Add Full-Screen Zoom Slides")
    print("=" * 70)
    print()

    # Paths
    pptx_path = "Inbound/MDF/Module_01_Foundation_Setup.pptx"
    images_dir = Path("Inbound/MDF/images")

    # Use timestamp to avoid file locking
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = f"Inbound/MDF/Module_01_With_Zoom_Slides_{timestamp}.pptx"

    # Image mappings (slide_index: image_filename)
    # Process in FORWARD order and track cumulative offset
    image_mappings = [
        (7, "slide_07_four_database_architecture.png", "Four-Database Architecture"),
        (12, "slide_12_four_warehouses.png", "Four Warehouses"),
        (18, "slide_18_role_hierarchy.png", "Role Hierarchy Pyramid"),
        (24, "slide_24_three_pillars.png", "Three Pillars Architecture"),
    ]

    # Check if presentation exists
    if not os.path.exists(pptx_path):
        print(f"Error: Presentation not found at {pptx_path}")
        print("Run: python generate_module01_slides.py first")
        return False

    # Check for images
    image_files = list(images_dir.glob("*.png"))
    if not image_files:
        print("No images found!")
        print("Generate images using Google Imagen first.")
        return False

    print(f"Loading presentation: {pptx_path}")
    prs = Presentation(pptx_path)
    original_slide_count = len(prs.slides)

    zoom_slides_added = 0
    offset = 0  # Track cumulative offset from inserted slides

    # Process images in FORWARD order, adjusting for previous insertions
    for original_idx, image_filename, description in image_mappings:
        image_path = images_dir / image_filename

        if image_path.exists():
            # Adjust index for previously inserted slides
            current_idx = original_idx + offset

            print(f"\nProcessing Slide {current_idx}: {description}")

            # Step 1: Remove watermark
            print(f"  [1/3] Removing watermark from image...")
            clean_image_path = remove_watermark_from_image(image_path, crop_percentage=10)

            # Step 2: Insert full-screen zoom slide RIGHT AFTER the current slide
            insert_position = current_idx + 1
            print(f"  [2/3] Inserting full-screen zoom slide at position {insert_position}...")
            zoom_slide = add_fullscreen_image_slide(prs, clean_image_path, insert_position)

            print(f"  [3/3] Zoom slide added - Original: Slide {current_idx}, Zoom: Slide {insert_position}")
            zoom_slides_added += 1

            # Increment offset for next iteration
            offset += 1

        else:
            print(f"\n[MISSING] {image_filename} - Skipping")

    # Save updated presentation
    if zoom_slides_added > 0:
        prs.save(output_path)
        print(f"\n{'=' * 70}")
        print(f"Success! Added {zoom_slides_added} full-screen zoom slides")
        print(f"{'=' * 70}")
        print(f"\nOriginal slides: {original_slide_count}")
        print(f"New slides: {len(prs.slides)}")
        print(f"Zoom slides added: {zoom_slides_added}")
        print(f"\nSaved to: {output_path}")
        print(f"\n{'=' * 70}")
        print("Presentation Structure:")
        print("=" * 70)
        print("Slide 7:  Four-Database Architecture (original with bullets)")
        print("Slide 8:  -> FULL-SCREEN ZOOM (no watermark)")
        print()
        print("Slide 13: Purpose-Specific Warehouses (original with bullets)")
        print("Slide 14: -> FULL-SCREEN ZOOM (no watermark)")
        print()
        print("Slide 20: Role Hierarchy (original with bullets)")
        print("Slide 21: -> FULL-SCREEN ZOOM (no watermark)")
        print()
        print("Slide 27: Three Pillars (original with bullets)")
        print("Slide 28: -> FULL-SCREEN ZOOM (no watermark)")
        print()
        print("=" * 70)
        print("Usage: Simply advance to next slide to show full-screen image")
        print("       Press ESC or click to return to normal flow")
        print("=" * 70)

        return True
    else:
        print("\nNo images found to add!")
        return False

if __name__ == "__main__":
    main()
