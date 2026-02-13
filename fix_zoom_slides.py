"""
Create presentation with zoom slides in correct order.
Copies all slides and inserts zoom slides at the right positions.
"""
import sys
import os
from pathlib import Path
from datetime import datetime

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "_video_automation"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import copy

def remove_watermark(image_path, crop_percent=10):
    """Crop bottom to remove watermark."""
    img = Image.open(image_path)
    width, height = img.size
    new_height = int(height * (100 - crop_percent) / 100)
    cropped = img.crop((0, 0, width, new_height))

    output_path = str(image_path).replace('.png', '_clean.png')
    cropped.save(output_path, 'PNG', quality=95)
    return output_path

def create_zoom_slide(prs, image_path):
    """Create a full-screen zoom slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Calculate dimensions for full-screen fit
    img = Image.open(image_path)
    img_w, img_h = img.size
    img_aspect = img_w / img_h

    slide_w, slide_h = 16.0, 9.0
    slide_aspect = slide_w / slide_h

    if img_aspect > slide_aspect:
        width = Inches(slide_w)
        height = Inches(slide_w / img_aspect)
        left = Inches(0)
        top = Inches((slide_h - (slide_w / img_aspect)) / 2)
    else:
        height = Inches(slide_h)
        width = Inches(slide_h * img_aspect)
        left = Inches((slide_w - (slide_h * img_aspect)) / 2)
        top = Inches(0)

    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    return slide

def copy_slide(source_slide, dest_prs):
    """Copy a slide to destination presentation."""
    # Get the source slide layout (use blank layout for destination)
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    # Copy background
    if source_slide.background.fill.type:
        new_slide.background.fill.solid()
        new_slide.background.fill.fore_color.rgb = source_slide.background.fill.fore_color.rgb

    # Copy all shapes
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide

def main():
    print("=" * 70)
    print("Creating Presentation with Zoom Slides in Correct Order")
    print("=" * 70)

    original_path = "Inbound/MDF/Module_01_Foundation_Setup.pptx"
    images_dir = Path("Inbound/MDF/images")

    # Slides that get zoom slides AFTER them (1-indexed)
    slides_with_zoom = {
        7: "slide_07_four_database_architecture.png",
        12: "slide_12_four_warehouses.png",
        18: "slide_18_role_hierarchy.png",
        24: "slide_24_three_pillars.png",
    }

    print(f"\nLoading original: {original_path}")
    original_prs = Presentation(original_path)

    print("Creating new presentation...")
    new_prs = Presentation()
    new_prs.slide_width = original_prs.slide_width
    new_prs.slide_height = original_prs.slide_height

    slide_count = 0
    zoom_count = 0

    # Copy each slide, inserting zoom slides where needed
    for idx in range(len(original_prs.slides)):
        slide_num = idx + 1  # 1-indexed

        # Copy original slide
        print(f"\nSlide {slide_num}: Copying original slide...")
        source_slide = original_prs.slides[idx]
        copy_slide(source_slide, new_prs)
        slide_count += 1

        # Check if this slide needs a zoom slide after it
        if slide_num in slides_with_zoom:
            image_file = slides_with_zoom[slide_num]
            image_path = images_dir / image_file

            if image_path.exists():
                zoom_position = slide_count + 1

                print(f"  -> Adding zoom slide at position {zoom_position}")
                print(f"     [1/2] Removing watermark...")
                clean_image = remove_watermark(image_path)

                print(f"     [2/2] Creating full-screen zoom...")
                create_zoom_slide(new_prs, clean_image)

                slide_count += 1
                zoom_count += 1

                print(f"     [OK] Zoom added - Original: {slide_num}, Zoom: {zoom_position}")

    # Save
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = f"Inbound/MDF/Module_01_Final_{timestamp}.pptx"

    print(f"\nSaving presentation...")
    new_prs.save(output_path)

    print(f"\n{'=' * 70}")
    print(f"SUCCESS!")
    print(f"={'=' * 70}")
    print(f"\nTotal slides: {slide_count}")
    print(f"  - Original slides: {len(original_prs.slides)}")
    print(f"  - Zoom slides: {zoom_count}")
    print(f"\nSaved to: {output_path}")

    print(f"\n{'=' * 70}")
    print("CORRECT SLIDE ORDER:")
    print("=" * 70)
    print("Slide 7:  Four-Database Architecture (with bullets)")
    print("Slide 8:  FULL-SCREEN ZOOM (no watermark)")
    print("Slide 9:  Why Separate Databases? (continues)")
    print()
    print("Slide 13: Purpose-Specific Warehouses (with bullets)")
    print("Slide 14: FULL-SCREEN ZOOM (no watermark)")
    print("Slide 15: Cost Attribution (continues)")
    print()
    print("Slide 20: Role Hierarchy (with bullets)")
    print("Slide 21: FULL-SCREEN ZOOM (no watermark)")
    print("Slide 22: Functional vs User Roles (continues)")
    print()
    print("Slide 27: Three Pillars (with bullets)")
    print("Slide 28: FULL-SCREEN ZOOM (no watermark)")
    print("Slide 29: Production Architecture (continues)")
    print(f"{'=' * 70}")

if __name__ == "__main__":
    main()
